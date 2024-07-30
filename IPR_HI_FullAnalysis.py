import pandas as pd
from sqlalchemy import create_engine
import textwrap as tw
import matplotlib.pyplot as plt
import seaborn as sns
import numpy as np
from scipy.stats import mannwhitneyu
from matplotlib.ticker import PercentFormatter
import matplotlib.dates as mdates
import scipy.stats.distributions as dist
import datetime as dt
import itertools
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.util import Inches, Pt
import time
import os
os.chdir('C:/Users/obriene/Projects/Health Inequalities Reporting')

#Get the current month and year for outputs
current_day = dt.datetime.today()
version_date = f'{current_day.strftime("%b")} {current_day.year}'
#Get date of last day of previous month for sql queries
end_date = (current_day.replace(day=1)
            - dt.timedelta(days=1))
start_date = end_date - relativedelta(years=1)
op_end_date = f'{end_date.strftime('%d-%B-%Y').upper()} 23:59:59'

# =============================================================================
# % Get data from queries
# ============================================================================= 
t0 = time.time()     
sdmart_engine = create_engine('mssql+pyodbc://@SDMartDataLive2/InfoDB?'\
                              'trusted_connection=yes&driver=ODBC+Driver+17'\
                              '+for+SQL+Server')
print('Reading in data...')
print('Clock Stops query running...')
# RTT Clock stops - Data Retrieval and formatting
query = f"""
DECLARE                       @StartDate AS DATETIME
DECLARE                       @Enddate AS DATETIME
-- SELECT data from a year ago to now
SET                           @Enddate = '{op_end_date}'
SET                           @StartDate = CAST(DATEADD(YEAR, -1, @Enddate) AS Date)
--Get RTT clock stops since 01/01/2022, disregarding those removed because they died
--NonAdmitted Clock Stops
SELECT rtt_nadm.compl_dttm, rtt_nadm.weeks_wait, rtt_nadm.days_wait,
ISNULL(REPLACE(rtt_nadm.nhs_number, ' ', ''), rtt_nadm.pasid) AS pat_no,
pat.pat_pcode, imd.[IndexValue] AS 'Decile ', Eth.[description],
spec.pfmgt_spec_desc AS specialty,
CASE WHEN specialty = 'Upper GI Surgery' THEN 'General, HpB, Oesophago-Gastric, Colorectal and Urology'
ELSE spec.slc_desc END AS SLC
--,case when alert.patnt_refno is not NULL then 'LD' else 'non-LD' end as 'LD Flag'
FROM infodb.dbo.rtt_daily_non_admitted_snapshot  rtt_nadm
--Get postcode from patients
LEFT JOIN PiMSMarts.dbo.patients pat
ON ISNULL(REPLACE(rtt_nadm.nhs_number, ' ', ''), rtt_nadm.pasid) = ISNULL(pat.nhs_number, pat.pasid)
--Get IMD from postcode
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
ON pat.pat_pcode = imd.PostcodeFormatted
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth
ON Eth.identifier = pat.ethgr
LEFT JOIN (SELECT DISTINCT pfmgt_spec, pfmgt_spec_desc, slc_desc FROM infoDB.dbo.vw_cset_specialties) spec
ON rtt_nadm.specialty = spec.pfmgt_spec
--Get LD patients from patient alerts
LEFT JOIN (SELECT DISTINCT patnt_refno
		   FROM PiMSMarts.dbo.Patient_Alert_Mart 
		   WHERE ODPCD_CODE = 'COM06' AND END_DTTM IS NULL) alert
ON pat.patnt_refno=alert.PATNT_REFNO
WHERE
(rtt_nadm.compl_dttm BETWEEN @StartDate AND @EndDate)--> '01/01/2022')
AND (rtt_nadm.provider IN ('RK900','89006','89999','NT200','NTY00') 
OR (rtt_nadm.provider ='XXXXX' 
AND rtt_nadm.clinic_code ='AC-S'))
AND cs_sorce_code <> 'DEATH' --Taken out those who were removed because they died
AND (cs_identifier <> '3' AND cs_sorce_code <> 'WLREM') -- if WLREM and 3, this also means patient died
AND imd.EndDate IS NULL
UNION
--Admitted Clock Stops
SELECT rtt_adm.admit_dttm AS compl_dttm, rtt_adm.weeks_wait_unadj, rtt_adm.days_wait_unadj,
ISNULL(rtt_adm.nhs_number,rtt_adm.pasid) AS pat_no, pat.pat_pcode, imd.[IndexValue] AS 'Decile ',
Eth.[description], specialty,
CASE WHEN specialty = 'Upper GI Surgery' THEN 'General, HpB, Oesophago-Gastric, Colorectal and Urology'
ELSE spec.slc_desc END AS SLC
--,case when alert.patnt_refno is not NULL then 'LD' else 'non-LD' end as 'LD Flag'
FROM InfoDB.dbo.vw_rtt_admitted_clock_stops  rtt_adm
--Get postcode from patients
LEFT JOIN PiMSMarts.dbo.patients pat
ON ISNULL(REPLACE(rtt_adm.nhs_number, ' ', ''), rtt_adm.pasid) = ISNULL(pat.nhs_number,pat.pasid)
--Get IMD from postcode
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
ON pat.pat_pcode = imd.PostcodeFormatted
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth
ON Eth.identifier = pat.ethgr
LEFT JOIN (SELECT DISTINCT pfmgt_spec_desc, slc_desc
		   FROM infoDB.dbo.vw_cset_specialties) spec
ON rtt_adm.specialty = spec.pfmgt_spec_desc
--Get LD patients from patient alerts
LEFT JOIN (SELECT DISTINCT patnt_refno
		   FROM PiMSMarts.dbo.Patient_Alert_Mart 
		   WHERE ODPCD_CODE = 'COM06' AND END_DTTM IS NULL) alert
ON pat.patnt_refno=alert.PATNT_REFNO
WHERE rtt_adm.type NOT IN ('Diag','NotTreated') --Removed Diag and Not Treated, so just Clock Stops
AND rtt_adm.admit_dttm BETWEEN @StartDate AND @EndDate--> '01/01/2022' 
AND rtt_adm.disch_outcome <> 'Died'
AND imd.EndDate IS NULL
"""
rtt_cs = pd.read_sql(query, sdmart_engine)
print('Clock Stops query complete')
#Get IMD deciles 1&2 and all others. Also don't include NaN days wait
rtt_cs['Decile '] = rtt_cs['Decile '].astype(float)

# RTT WL - Data Retrieval
print('RTT Wait List query running...')
wl_query = """
--Will return current RTT incomplete waiting list position
SELECT imd.[IndexValue] AS 'Decile ', RTT.current_LOW, eth.[description],
Specialty_Referred_to AS specialty, spec.slc_desc AS SLC
FROM [InfoDB].[dbo].[rtt_daily_incomplete_pathways_snapshot] RTT
LEFT JOIN PiMSMarts.dbo.referrals Ref
ON RTT.refrl_refno = ref.refrl_refno 
LEFT JOIN PiMSMarts.dbo.patients pat
ON Ref.patnt_refno = pat.patnt_refno
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth
ON Eth.identifier = pat.ethgr
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
ON pat.pat_pcode = imd.PostcodeFormatted
--left join [InfoDB].[dbo].[IMDScore] imd
			--on pat.pat_pcode = imd.PCD2
LEFT JOIN (SELECT DISTINCT pfmgt_spec, pfmgt_spec_desc, slc_desc
		   FROM infoDB.dbo.vw_cset_specialties) spec
		   ON RTT.pfmgt_spec = spec.pfmgt_spec
WHERE run_date = (SELECT MAX(run_date)
FROM [InfoDB].[dbo].[rtt_daily_incomplete_pathways_snapshot])
AND imd.EndDate IS NULL
"""
rtt_incomp = pd.read_sql(wl_query, sdmart_engine)
print('RTT Wait List query complete') 
#Make Decile column numeric
rtt_incomp['Decile '] = rtt_incomp['Decile '].astype(float)

# Non-F2F Section Data retrieval
print('OP query running...')
op_query = f"""
DECLARE                       @dtmStartDate AS DATETIME
DECLARE                       @dtmEnddate AS DATETIME
-- select data from a year ago to now
SET                           @dtmEnddate = '{op_end_date}'
SET                           @dtmStartDate = CAST(DATEADD(YEAR, -1, @dtmEnddate) AS Date)

SELECT opact.pasid,
Visit = CASE WHEN opact.visit IN ('1','2') THEN 'F2F'
		WHEN opact.visit IN ('3','4') THEN 'Non-F2F' END,
DATEDIFF(YEAR, pats.pat_dob,opact.start_dttm) AS Age,
CASE WHEN DATEDIFF(YEAR, pats.pat_dob,opact.start_dttm) < 20 THEN '0-19'
	 WHEN DATEDIFF(YEAR,pats.pat_dob,opact.start_dttm) BETWEEN 20 AND 29 THEN '20-29'
	 WHEN DATEDIFF(YEAR,pats.pat_dob,opact.start_dttm) BETWEEN 30 AND 39 THEN '30-39'
	 WHEN DATEDIFF(YEAR,pats.pat_dob,opact.start_dttm) BETWEEN 40 AND 49 THEN '40-49'
	 WHEN DATEDIFF(YEAR,pats.pat_dob,opact.start_dttm) BETWEEN 50 AND 59 THEN '50-59'
	 WHEN DATEDIFF(YEAR,pats.pat_dob,opact.start_dttm) BETWEEN 60 AND 69 THEN '60-69'
	 WHEN DATEDIFF(YEAR,pats.pat_dob,opact.start_dttm) BETWEEN 70 AND 79 THEN '70-79'
	 WHEN DATEDIFF(YEAR,pats.pat_dob,opact.start_dttm)  >= 80 THEN '80+'
	 END AS Age_range,
specialty = spec.pfmgt_spec_desc, slc = spec.slc_desc, [Eth].[description] AS Ethnicity,
imd.[IndexValue], pats.disabled_yn                
FROM Pimsmarts.dbo.outpatients opact
LEFT JOIN InfoDB.dbo.vw_cset_specialties spec
ON opact.local_spec = spec.local_spec 
LEFT JOIN PiMSMarts.dbo.patients pats
ON opact.patnt_refno = pats.patnt_refno
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth 
ON Eth.identifier = pats.ethgr
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
ON opact.pat_pcode = imd.PostcodeFormatted
WHERE opact.cancr_dttm IS NULL 
AND opact.start_dttm BETWEEN @dtmStartDate AND @dtmEnddate
--these are standard exclusions
AND (opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN ('RK900','89006','89999','NT200','NTY00') OR (opact.provider ='XXXXX'
AND opact.clinic_code = 'AC-S'))
AND (opact.pat_surname NOT LIKE 'ZZ%' AND opact.pat_surname NOT LIKE 'XX%')
AND opact.sctyp='OTPAT' AND opact.session_code IS NOT NULL
AND spec.pfmgt_spec <> 'ZZ' -- remove non UHP activity
AND imd.EndDate IS NULL
"""
op_data = pd.read_sql(op_query, sdmart_engine)
print('OP query complete')
op_data['IndexValue'] = op_data['IndexValue'].astype(float)
#print timings
t1 = time.time()
print(f'Queries run in {(t1-t0)/60} mins')
# =============================================================================
# % Functions
# =============================================================================
#Function to plot labels on bars
def show_values_on_bars(axs,percentage = True, rounded = 2, numbers = None):
    def _show_on_single_plot(ax):
        counter = 0       
        for p in ax.patches:
            if p._height !=0:
                if percentage:
                    _x = p.get_x() + p.get_width() / 2
                    if (p.xy[1] == 0) and (p._height < 1):#p._height < 0.50:
                        _y = p.get_y() + p.get_height() + 0.01

                        if rounded == 2:
                            value = '{:.2f}'.format(p.get_height()*100)
                        elif rounded == 0:
                            value = '{:.0f}'.format(p.get_height()*100)
                        if numbers:
                            ax.text(_x, _y+0.02, value+"%\n("+str(numbers[counter])+")",
                                    ha="center")
                            counter = counter + 1
                        else:
                            ax.text(_x, _y, value+"%", ha="center")
                else:
                    _x = p.get_x() + p.get_width() / 2
                    _y = p.get_y() + p.get_height()+0.3
                    value = '{:.1f}'.format(p.get_height())
                    ax.text(_x, _y, value, ha="center") 

    if isinstance(axs, np.ndarray):
        for idx, ax in np.ndenumerate(axs):
            _show_on_single_plot(ax)
    else:
        _show_on_single_plot(axs)
        
#Function to show totals at top of bars
def show_totals(axs, totals):
    def _show_single_totals(ax):
        counter = 0
        for p in ax.patches:
            if p._height != 0:
                _x = p.get_x() + p.get_width() / 2
                if ((p.xy[1] > 0) and (p._height < 1)) or (p._height == 1):#p._height > 0.50:
                    _y = p.get_y() + p.get_height()-0.05
                    ax.text(_x, _y, str(totals[counter]), ha="center")
                    counter = counter + 1
    if isinstance(axs, np.ndarray):
        for idx, ax in np.ndenumerate(axs):
            _show_single_totals(ax)
    else:
        _show_single_totals(axs)
        
#Pie formatting
def autopct_format(values):
    def my_format(pct):
        total = sum(values)
        val = int(round(pct * total / 100.0))
        return '{:.1f}%\n({v:d})'.format(pct, v=val)
    return my_format

#Significance bars
def label_diff(ax, i, j, text, X, Y):
    x = (X[i] + X[j]) / 2
    y = 1.07 * max(Y[i], Y[j])
    props = {'connectionstyle':'bar', 'arrowstyle':'-', 'shrinkA':20,
             'shrinkB':20,'linewidth':2}
    ylims = ax.get_ylim()[1] - ax.get_ylim()[0]
    #If its a percentage plot, don't need extra y increase
    if max(Y) == min(Y) == 1:
        ax.annotate(text, xy=(x, y*1.05), zorder=10, ha='center',
                    annotation_clip=False)
        ax.annotate('', xy=(X[i], y*0.87), xytext=(X[j], y*0.87),
                    arrowprops=props, annotation_clip=False)
    else:
        ax.annotate(text, xy=(x, y+0.2*ylims), zorder=10, ha='center')
        ax.annotate('', xy=(X[i], y), xytext=(X[j], y), arrowprops=props)        

#Test whether two proportions are different
def propHypothesisTest(p1, p2, n1, n2, alpha = 0.05):
    #Following:https://medium.com/analytics-vidhya/testing-a-difference-in-population-proportions-in-python-89d57a06254
    #p1 and p2 are the proportions of each dataset falling in the 'yes' category
    #n1 and n2 are the total number of datapoints in each dataset
    #Alpha is the signifficance threshold (10% for this 2-tailed test)
    #Null Hypothesis: Proportions equal
    #Alternative: Proportions significantly different
    #First, find the standard error
    #For this, we need the total proportion with a yes classification
    p = (n1*p1 + n2*p2)/(n1 + n2)
    se = np.sqrt(p*(1-p)*((1/n1) + (1/n2)))
    #Next, calculate the test statistic:
        #(best estimate - hypothesized estimate)/standard error
        #best estimate = p1-p2, hypothesized = 0(as p1 and p2 are equal)
    if se == 0:
        return None
    test_stat = (p1-p2)/se 
    #This gives number of standard deviations from hypothesized estimate
    #From the test statistic, get the p-value
    pvalue = 2 * dist.norm.cdf(-np.abs(test_stat)) # Multiplied by two indicates a two tailed testing.
    return pvalue

#function to get list of counts to account for missing values
def counts_list(counts, options):
    c1 = (counts[options[0]] if options[0] in counts.index else 0)
    c2 = (counts[options[1]] if options[1] in counts.index else 0)
    return [c1, c2]

# =============================================================================
#     #Initial analysis and Sort data
# =============================================================================
# % Analyse data as a whole before going to specialty level
# % RTT Clock Stops Analysis Section
#Get cs for given specialty
rtt_slc = rtt_cs.copy()
#Get IMD deciles 1&2 and all others. Also don't include NaN days wait
imd_1_2 = rtt_slc.loc[(rtt_slc['Decile '].isin([1,2]))
                      & (~pd.isnull(rtt_slc['days_wait']))].copy()
imd_3_10 = rtt_slc.loc[(~rtt_slc['Decile '].isin([1,2]))
                       & (~pd.isnull(rtt_slc['days_wait']))
                       & (~pd.isnull(rtt_slc['Decile ']))].copy()
#Get minority ethnicities and non
me = rtt_slc.loc[(~rtt_slc['description'].isin(
                  ['Unknown', 'Unwilling to answer', 'White British']))
                 & (~pd.isnull(rtt_slc['days_wait']))].copy()
wb = rtt_slc.loc[(rtt_slc['description'] == 'White British')
                 & (~pd.isnull(rtt_slc['days_wait']))].copy()
#Find median length of wait for each of these populations
me_med_low = me['days_wait'].median()
wb_med_low = wb['days_wait'].median()
imd_1_2_med_low = imd_1_2['days_wait'].median()
imd_3_10_med_low = imd_3_10['days_wait'].median()

# =============================================================================
#     #Make a bar chart of the median LoW for IMD and Eth categories SLIDE 2
# =============================================================================
#Test for statistical difference in wait length for IMD and Ethnicity
pval_rtt_cs_IMD = mannwhitneyu(imd_1_2['days_wait'].tolist(),
                               imd_3_10['days_wait'].tolist(),
                               alternative = 'greater')[1]
pval_rtt_cs_eth = mannwhitneyu(me['days_wait'].tolist(),
                               wb['days_wait'].tolist(),
                               alternative = 'greater')[1]
#Plot bar chart
fig,ax = plt.subplots(1,1)
ax.bar([1, 2, 3, 4],
       [me_med_low, wb_med_low, imd_1_2_med_low, imd_3_10_med_low ],
       color = ['royalblue','lightskyblue','seagreen','lightgreen'],
       edgecolor='black')
ax.set_xticks([1,2,3,4])
ax.set_xticklabels(['Ethnic\n Minority\n('+f"{me.shape[0]:,.0f}"+' clock \nstops)',
                'White British\n('+f"{wb.shape[0]:,.0f}"+' \nclock stops)',
                'IMD 1-2\n('+f"{imd_1_2.shape[0]:,.0f}"+'\n clock stops)',
                'IMD 3-10\n('+f"{imd_3_10.shape[0]:,.0f}"+'\n clock stops)'])
ax.set_ylabel('Median LoW (Days)')
ax.set_title('Median Length of Wait for RTT')
show_values_on_bars(ax, percentage = False)

#Add labels for significant difference if they occur.
if pval_rtt_cs_eth < 0.05:
    label_diff(ax, 0, 1, 'Significantly longer\n LoW for EM patients',
               [1,2,3,4],
               [me_med_low,wb_med_low,imd_1_2_med_low,imd_3_10_med_low])
else:
     label_diff(ax, 0, 1, 'No significant difference',
                [1,2,3,4],
                [me_med_low,wb_med_low,imd_1_2_med_low,imd_3_10_med_low])
     
if pval_rtt_cs_IMD < 0.05:
     label_diff(ax, 2, 3, 'Significantly longer\n LoW for IMD 1&2 patients',
                [1,2,3,4],
                [me_med_low,wb_med_low,imd_1_2_med_low,imd_3_10_med_low])
else:
      label_diff(ax, 2, 3, 'No significant difference',
                 [1,2,3,4],
                 [me_med_low,wb_med_low,imd_1_2_med_low,imd_3_10_med_low])
plt.ylim(ymax=ax.get_ylim()[1]*1.4)
plt.savefig('plots/Slide 2.png', bbox_inches='tight')

# =============================================================================
# # ETHNICITY SLC RTT MEDIAN LOW - statistical tests SLIDE 3
# =============================================================================
slc_unique = rtt_cs['SLC'].unique()
RTT_LOW_eth_pvals = []
wb_med_low_spec = []
me_med_low_spec = []
slc_unique = rtt_cs['SLC'].unique()

for spec in slc_unique:
    wb_spec = wb.loc[wb['SLC'] == spec]
    me_spec = me.loc[me['SLC'] == spec]
    #Only perform Mood's median test if there are t least 15 samples in each 
    #pop. Also, if both medians are zero, median test cannot work properly, so 
    #exclude
    if ((wb_spec.shape[0] >= 15) and (me_spec.shape[0] >= 15)
        and (wb_spec['days_wait'].median() != 0)
        and (me_spec['days_wait'].median() != 0)):
        RTT_LOW_eth_pvals.append(mannwhitneyu(wb_spec['days_wait'].tolist(),
                                              me_spec['days_wait'].tolist())[1])
    else:
        RTT_LOW_eth_pvals.append(np.nan)
    #add median length of stay for each slc
    wb_med_low_spec.append(wb_spec['days_wait'].median())
    me_med_low_spec.append(me_spec['days_wait'].median())

#Make a df of the results
RTT_LOW_eth_pvals = pd.DataFrame({'SLC':slc_unique, 
                                  'p-value':RTT_LOW_eth_pvals,
                                  'Median LoW WB':wb_med_low_spec,
                                  'Median LoW ME':me_med_low_spec})
RTT_LOW_eth_pvals['WB-ME'] = (RTT_LOW_eth_pvals['Median LoW WB']
                              - RTT_LOW_eth_pvals['Median LoW ME'])

#Filter to statistically significant
rtt_low_eth_slc_plt = pd.melt(RTT_LOW_eth_pvals.loc[RTT_LOW_eth_pvals['p-value']
                                                    <0.025],
                              id_vars=['SLC'],
                              value_vars = ['Median LoW ME', 'Median LoW WB'])

#Plot
fig, ax_spec_eth = plt.subplots(1, 1, figsize = (9,4))
sns.barplot(data=rtt_low_eth_slc_plt, x='SLC', hue = 'variable', y='value',
            ax=ax_spec_eth, palette=['royalblue','lightskyblue'])
legend = ax_spec_eth.get_legend()
# #Get seaborn legend
handles = legend.legend_handles
ax_spec_eth.legend(handles, ['Ethnic Minority', 'White British'],
                   bbox_to_anchor = (1,1))
ax_spec_eth.set_ylabel('Median LoW (days)')
ax_spec_eth.set_xticks(ax_spec_eth.get_xticks())
labels = [tw.fill(l, 20) for l in rtt_low_eth_slc_plt.SLC.unique()]
ax_spec_eth.set_xticklabels(labels=labels, fontsize=8)
plt.savefig('plots/Slide 3.png', bbox_inches='tight')

# =============================================================================
# #IMD RTT LoW by SLC SLIDE 4
# =============================================================================
slc_unique = rtt_slc['SLC'].unique()
RTT_LOW_IMDSLC_pvals = []
imd_1_2_med_low_slc = []
imd_3_10_med_low_slc = []

for slc in slc_unique:
    imd_1_2_spec = imd_1_2.loc[imd_1_2['SLC'] == slc].copy()
    imd_3_10_spec = imd_3_10.loc[imd_3_10['SLC'] == slc].copy()
    #Only perform Mood's median test if there are at least 15 samples in each 
    # population. Also, if both medians are zero,
    # median test cannot work properly,so exclude
    if ((imd_1_2_spec.shape[0] >= 15) and (imd_3_10_spec.shape[0] >= 15)
        and (imd_1_2_spec['days_wait'].median() != 0)
        and (imd_3_10_spec['days_wait'].median() != 0)):
        RTT_LOW_IMDSLC_pvals.append(
            mannwhitneyu(imd_1_2_spec['days_wait'].tolist(),
                         imd_3_10_spec['days_wait'].tolist())[1])
    else:
        RTT_LOW_IMDSLC_pvals.append(np.nan)
    #add median length of stay for each slc
    imd_1_2_med_low_slc.append(imd_1_2_spec['days_wait'].median())
    imd_3_10_med_low_slc.append(imd_3_10_spec['days_wait'].median())

#Make a df of the results
RTT_LOW_IMDSLC_pvals = pd.DataFrame({'SLC':slc_unique,
                                    'p-value':RTT_LOW_IMDSLC_pvals,
                                    'Median LoW IMD 1-2':imd_1_2_med_low_slc,
                                    'Median LoW IMD 3-10':imd_3_10_med_low_slc})
RTT_LOW_IMDSLC_pvals['IMD1_2-IMD3_10'] = (RTT_LOW_IMDSLC_pvals['Median LoW IMD 1-2']
                                          - RTT_LOW_IMDSLC_pvals['Median LoW IMD 3-10'])
#Filter to statistically significant
rtt_low_imd_slc_plt = pd.melt(RTT_LOW_IMDSLC_pvals.loc[
    RTT_LOW_IMDSLC_pvals['p-value'] < 0.025], id_vars=['SLC'],
    value_vars = ['Median LoW IMD 1-2', 'Median LoW IMD 3-10'])
#plot
fig, ax_spec_imd = plt.subplots(1,1,figsize = (8,4))
sns.barplot(data=rtt_low_imd_slc_plt, x='SLC', hue='variable', y='value',
            ax=ax_spec_imd, palette=['seagreen','lightgreen'])
legend = ax_spec_imd.get_legend()
handles = legend.legend_handles
ax_spec_imd.legend(handles, ['IMD 1-2','IMD 3-10'],
               bbox_to_anchor = (1,1))
ax_spec_imd.set_ylabel('Median LoW (days)')
ax_spec_imd.set_xticks(ax_spec_imd.get_xticks())
labels = [tw.fill(l, 21) for l in rtt_low_imd_slc_plt.SLC.unique()]
ax_spec_imd.set_xticklabels(labels = labels)
plt.savefig('plots/Slide 4.png', bbox_inches='tight')

# % RTT Incomplete WL Analysis
# ========================================================================
#     #>52 week wait analysis SLIDE 5
# ========================================================================
#############All week waits
#####IMD
#Get a version with no NaNs for deciles
rtt_incomp_dec = rtt_incomp.loc[~pd.isnull(rtt_incomp['Decile '])].copy()
#Make new columns with IMD1-2 or IMD3-10
rtt_incomp_dec['value'] = np.where(rtt_incomp_dec['Decile '].isin([1, 2]),
                                   'IMD 1-2', 'IMD 3-10')
rtt_incomp_dec['type'] = 'IMD'
#####Ethnicity
#Get a version without unknown ethnicities
rtt_incomp_eth = rtt_incomp.loc[~rtt_incomp['description']
                                .isin(['Unknown','Unwilling to answer'])].copy()
#Add column for white british and ethnic minority split
rtt_incomp_eth['value'] = np.where(rtt_incomp_eth['description']
                                   == 'White British',
                                   'White British', 'Ethnic Minority')
rtt_incomp_eth['type'] = 'Ethnicity'

############# > 52 week wait
#####IMD
#Filter to >25 weeks and remove decile nans
rtt_incomp_52_dec = rtt_incomp.loc[(rtt_incomp['current_LOW'] > 364)
                                   & (~pd.isnull(rtt_incomp['Decile ']))].copy()
#Make columns for IMD and ethnicity as above
rtt_incomp_52_dec['value'] = np.where(rtt_incomp_52_dec['Decile '].isin([1, 2]),
                                   'IMD 1-2', 'IMD 3-10')
rtt_incomp_52_dec['type'] = 'IMD\n (>52 Week Wait)'
#####Ethnicity
#Remove unknown ethnicities
rtt_incomp_52_eth = rtt_incomp.loc[(rtt_incomp['current_LOW'] > 364)
                                   & (~rtt_incomp['description'].isin(
                                    ['Unknown', 'Unwilling to answer']))].copy()
#Add column for white british and ethnic minority split
rtt_incomp_52_eth['value'] = np.where(rtt_incomp_52_eth['description']
                                      == 'White British',
                                   'White British', 'Ethnic Minority')
rtt_incomp_52_eth['type'] = 'Ethnicity \n(>52 Week Wait)'

##############Hypothesis testing
#####IMD
#Get counts
total_num_IMD = rtt_incomp_dec.shape[0]
total_num_IMD_52 = rtt_incomp_52_dec.shape[0]
n_IMD12 = rtt_incomp_dec[rtt_incomp_dec['value'] == 'IMD 1-2'].shape[0]
n_IMD12_52 = rtt_incomp_52_dec[rtt_incomp_52_dec['value'] == 'IMD 1-2'].shape[0]
#test hypothesis
pval_rtt_incomp_imd = propHypothesisTest((n_IMD12 / total_num_IMD),
                                         (n_IMD12_52 / total_num_IMD_52),
                                          n_IMD12, n_IMD12_52, alpha=0.05)
imd_str = 'a' if pval_rtt_incomp_imd < 0.05 else 'no'

#####Ethnicity
#get counts
total_num_eth = rtt_incomp_eth.shape[0]
total_num_eth_52 = rtt_incomp_52_eth.shape[0]
n_em = rtt_incomp_eth[rtt_incomp_eth['value'] == 'Ethnic Minority'].shape[0]
n_em_52 = rtt_incomp_52_eth[rtt_incomp_52_eth['value']
                            == 'Ethnic Minority'].shape[0]
#test hypothesis
pval_rtt_incomp_eth = propHypothesisTest((n_em / total_num_eth),
                                         (n_em_52 / total_num_eth_52),
                                          n_em, n_em_52, alpha=0.05)
eth_str = 'a' if pval_rtt_incomp_eth < 0.05 else 'no'

#####Plot
#Concat all dataframes into one large one
rtt_in_full = pd.concat([rtt_incomp_dec, rtt_incomp_eth,
                         rtt_incomp_52_dec, rtt_incomp_52_eth],
                         ignore_index=True)  
#Plot a filled bar
fig, ax_52WW = plt.subplots(1,1)
hue_order = ['White British', 'Ethnic Minority', 'IMD 3-10', 'IMD 1-2']
sns.histplot(data=rtt_in_full.sort_values(['type']), x='type', hue='value',
             multiple='fill', shrink=0.6, hue_order=hue_order,
             palette=['lightskyblue','royalblue','lightgreen','seagreen'],
             alpha=1, ax=ax_52WW)
ax_52WW.yaxis.set_major_formatter(PercentFormatter(1))
ax_52WW.set_xlabel('')
ax_52WW.set_ylabel('Percentage of Patients')
ax_52WW.set_title('Patients on RTT Waiting List')
legend = ax_52WW.get_legend()
# #Get seaborn legend
handles = legend.legend_handles
ax_52WW.legend(handles,
               ['White British', 'Ethnic Minority', 'IMD 3-10', 'IMD 1-2'],
               bbox_to_anchor=(1,1))
#Make a list of the numbers to include under the percentages
numbers = [n_IMD12, n_IMD12_52, n_em, n_em_52]
totals = [total_num_IMD, total_num_IMD_52, total_num_eth, total_num_eth_52]
show_values_on_bars(ax_52WW, numbers = [i for i in numbers if i!=0])
show_totals(ax_52WW, totals)
#add text boxes
fig.text(.95, .4,
         tw.fill(f"There is {eth_str} significant difference between the "\
                 "proportion of ethnic minority patients on the RTT waiting "\
                 "list and those waiting for > 52 weeks.", 40),
         ha='center', clip_on=False, fontsize=10,
         bbox=dict(boxstyle='round,pad=0.5', fc='none', ec='black'))
fig.text(.95, .2,
         tw.fill(f"There is {imd_str} significant difference between the "\
                 "proportion of IMD 1&2 patients on the RTT waiting list and "\
                 "those waiting for > 52 weeks.", 40),
         ha='center', clip_on=False, fontsize=10,
         bbox=dict(boxstyle='round,pad=0.5', fc='none', ec='black'))
plt.tight_layout()
plt.savefig('plots/Slide 5.png', bbox_inches='tight')

# ========================================================================
#     #Digital Access Analysis
# ========================================================================
#Set title text
title_text = 'Outpatient Appointments'
hue_order = ['F2F','Non-F2F']

# ======================================================================
#         #IMD (SLIDES 7-9)
# ======================================================================
#########Decile (2 category)  SLIDE 7
#Remove missing data
op_data_temp_imd = op_data.loc[~pd.isnull(op_data['IndexValue'])].copy()
op_data_temp_imd['value'] = np.where(op_data_temp_imd['IndexValue']
                                            .isin([1, 2]), 'IMD 1-2',
                                            'IMD 3-10')
#Bar plots
fig1, ax_op_dec2 = plt.subplots(1, 1, figsize = (4,4))
sns.histplot(data=op_data_temp_imd.sort_values('value'), x='value',
                hue="Visit", hue_order=hue_order, multiple="fill",
                discrete=True,  palette=['lemonchiffon','orange'], alpha=1,
                ax=ax_op_dec2)
#Make a list of the numbers to include under the percentages
nonf2f_counts = op_data_temp_imd.loc[op_data_temp_imd['Visit'] == 'Non-F2F',
                                    'value'].value_counts()
num_list = counts_list(nonf2f_counts, ['IMD 1-2', 'IMD 3-10'])
counts = op_data_temp_imd['value'].value_counts()
tot_list = counts_list(counts, ['IMD 1-2', 'IMD 3-10'])
show_values_on_bars(ax_op_dec2, rounded=0,
                    numbers = [j for j in num_list if j!=0])
show_totals(ax_op_dec2, tot_list)
ax_op_dec2.yaxis.set_major_formatter(PercentFormatter(1))
ax_op_dec2.set_xlabel('')
ax_op_dec2.set_ylabel('Percentage of ' + title_text)
legend = ax_op_dec2.get_legend()
# #Get seaborn legend
handles = legend.legend_handles
ax_op_dec2.legend(handles, hue_order, title='Appointment Type',
                    bbox_to_anchor=(1,1))
#Proportion testing
pval_op_imd = propHypothesisTest(num_list[0]/tot_list[0],
                                    num_list[1]/tot_list[1],
                                    num_list[0], num_list[1], alpha=0.025)
#Add bars to show if there is a significant difference
slide7_text = ('Significant difference' if pval_op_imd < 0.025
        else ' No significant difference')
label_diff(ax_op_dec2, 0, 1, slide7_text, [0,1], [1,1])
plt.savefig('plots/' + title_text + ' Slide 7.png', bbox_inches='tight')

# ======================================================================
#########IMD 1&2 vs 3-10 non-f2f plots and testing  SLIDE 8
spec_unique = op_data['specialty'].unique()
#Make an empty list to store pvalues in
spec_imd2_pvals = []
spec_12_nf2f_perc = []
spec_310_nf2f_perc = []
for spec in spec_unique:
    #Separate by EM
    op_data_temp = op_data_temp_imd.loc[op_data_temp_imd['specialty']
                                        == spec].copy()
    #Make a list of the numbers to include under the percentages
    #could not shorten as sometimes there isn't data for these.
    numbers = [op_data_temp[(op_data_temp['value'] == 'IMD 1-2')
                            & (op_data_temp['Visit'] == 'Non-F2F')]
                            .shape[0],
                op_data_temp[(op_data_temp['value'] == 'IMD 3-10')
                            & (op_data_temp['Visit'] == 'Non-F2F')]
                            .shape[0]]
    totals = [op_data_temp[(op_data_temp['value'] == 'IMD 1-2')].shape[0],
                op_data_temp[(op_data_temp['value'] == 'IMD 3-10')].shape[0]]
    
    if totals[0] == 0:
        spec_12_nf2f_perc.append(0)
    else:
        spec_12_nf2f_perc.append(100*(numbers[0]/totals[0]))
    if totals[1] == 0:
        spec_310_nf2f_perc.append(0)
    else:
        spec_310_nf2f_perc.append(100*(numbers[1]/totals[1]))
    #Hypothesis test for EM Non-F2F
    #If no non-F2F patients are seen in either or both ethnicity categories, 
    #cannot perform hypothesis test
    if ((all(j > 50 for j in numbers))
    or ((all(j > 50 for j in totals)) and (all(j > 0 for j in numbers)))):
        spec_imd2_pvals.append(propHypothesisTest(numbers[0]/totals[0],
                                                    numbers[1]/totals[1],
                                                    numbers[0], numbers[1],
                                                    alpha=0.025))
    else:
        spec_imd2_pvals.append(None)
#Create dataframe
spec_imd2_pvals = pd.DataFrame({'Specialty':spec_unique,
                                'p-value':spec_imd2_pvals,
                                'IMD 1&2 Non-F2F Perc':spec_12_nf2f_perc,
                                'IMD3-10 Non-F2F Perc':spec_310_nf2f_perc})
#Barplot
nf2f_imd2_spec_plt = pd.melt(spec_imd2_pvals[spec_imd2_pvals['p-value']
                                                <0.025],
                                id_vars=['Specialty'],
                                value_vars=['IMD 1&2 Non-F2F Perc',
                                            'IMD3-10 Non-F2F Perc'])
#Sort op_data by specialty
nf2f_imd2_spec_plt.sort_values('Specialty', inplace=True)
#Order this op_data by the differences between specs. First find the differences
differences = [item for item in (nf2f_imd2_spec_plt['value'].diff()
                                    .iloc[1::2].tolist())
                for j in range(2)]

#Add as new column to op_data
nf2f_imd2_spec_plt['diffs'] = differences
fig, ax_nf2fspec_imd2 = plt.subplots(1, 1, figsize=(9,4))
sns.barplot(data=nf2f_imd2_spec_plt.sort_values('diffs'), x='Specialty',
            hue='variable', y='value', ax=ax_nf2fspec_imd2,
            palette=['seagreen','lightgreen'])
legend = ax_nf2fspec_imd2.get_legend()
# #Get seaborn legend
handles = legend.legend_handles
ax_nf2fspec_imd2.legend(handles, ['IMD 1&2', 'IMD 3-10'],
                        bbox_to_anchor=(1,1))
if title_text == 'Outpatients':
    ax_nf2fspec_imd2.set_ylabel(tw.fill('Percentage of ' + title_text
                                        + ' seen non-F2F', 40))
else:
    ax_nf2fspec_imd2.set_ylabel(tw.fill('Percentage of ' + title_text
                                        + ' carried out non-F2F',40))
ax_nf2fspec_imd2.set_xticks(ax_nf2fspec_imd2.get_xticks())
labels = [tw.fill(l, 16) for l
            in nf2f_imd2_spec_plt.sort_values('diffs')['Specialty'].unique()]
ax_nf2fspec_imd2.set_xticklabels(labels = labels)
plt.tight_layout()
plt.savefig('plots/' + title_text + ' Slide 8.png', bbox_inches='tight')

# ======================================================================
############Decile (10 category) SLIDE 9
#Make a plot of the proportions
fig1, ax_op_dec = plt.subplots(1, 1, figsize = (7,5))
sns.histplot(data=op_data_temp_imd, x='IndexValue', hue='Visit',
                hue_order=hue_order, multiple='fill', discrete=True, 
                palette=['lemonchiffon','orange'], alpha=1, ax=ax_op_dec)
#Make a list of the numbers to include under the percentages
num_list = []
tot_list = []
for ii in range(1,11):
    num_list.append(op_data_temp_imd[(op_data_temp_imd['IndexValue'] == ii)
                                    & (op_data_temp_imd['Visit'] == 'Non-F2F')
                                    ].shape[0])
    tot_list.append(op_data_temp_imd[(op_data_temp_imd['IndexValue'] == ii)
                                    ].shape[0])

show_values_on_bars(ax_op_dec, rounded=0,
                    numbers=[ii for ii in num_list if ii!=0])
show_totals(ax_op_dec, tot_list)
#show_values_on_bars(ax_op_dec,rounded=0)
ax_op_dec.yaxis.set_major_formatter(PercentFormatter(1))
ax_op_dec.set_xlabel('IMD Decile')
ax_op_dec.set_ylabel('Percentage of ' + title_text)
legend = ax_op_dec.get_legend()
# #Get seaborn legend
handles = legend.legend_handles
ax_op_dec.legend(handles, hue_order,
                    title = 'Appointment Type',
                    bbox_to_anchor = (1,1))
plt.savefig('plots/' + title_text + ' Slide 9.png', bbox_inches='tight')

#Get all pairs of decile numbers with no repeats and hypothesis test
imd_pvals = []
pair_list = itertools.combinations(range(0,10),2)
for r in pair_list:
    pvalue = propHypothesisTest(num_list[r[0]]/tot_list[r[0]],
                                num_list[r[1]]/tot_list[r[1]],
                                num_list[r[0]], num_list[r[1]])
    #Record statistically significant results
    if pvalue < 0.025:
        imd_pvals.append([pvalue, r])

#Make a op_data of the results
imd_prop_test = pd.DataFrame(imd_pvals, columns=['pvals', 'pairs'])
imd_list = [j for j in range(1,11)]
imd_prop_test['imd1'] = (imd_prop_test['pairs']
                            .apply(lambda x: imd_list[x[0]]))
imd_prop_test['imd2'] = (imd_prop_test['pairs']
                            .apply(lambda x: imd_list[x[1]]))
#write a string of the IMD pairs that have a significant difference
pairs = (imd_prop_test.groupby('imd1', as_index=False)['imd2']
            .apply(list).values.tolist())
slide9_text = ''
for pair in pairs:
    imd = pair[0]
    lst = pair[1]
    seccond_str = ('all other IMD values' if len(lst) == (10 - imd)
                    else ('IMDs ' +   ', '.join(map(str, lst))))
    string = f'IMD {imd} and {seccond_str}\n'
    slide9_text += string

# ====================================================================
#         #Ethnicity  SLIDES 10-11
# ====================================================================
#############Ethnicity bar plot SLIDE 10
#Get rid of missing values
op_data_temp_eth = op_data.loc[~op_data['Ethnicity']
                        .isin(['Unknown', 'Unwilling to answer'])].copy()
op_data_temp_eth['value'] = np.where(op_data_temp_eth['Ethnicity']
                                        == 'White British',
                                        'White British', 'Ethnic Minority')
#Make a plot of the proportions 
fig, ax_op_eth = plt.subplots(1, 1, figsize = (4,4))
sns.histplot(data=op_data_temp_eth.sort_values('value'), x='value',
                hue="Visit", hue_order=hue_order, multiple="fill",
                discrete=True, palette=['lemonchiffon','orange'], alpha=1,
                ax=ax_op_eth)
ax_op_eth.yaxis.set_major_formatter(PercentFormatter(1))
ax_op_eth.set_xlabel('')
#Make a list of the numbers to include under the percentages
nonf2f_counts = op_data_temp_eth.loc[op_data_temp_eth['Visit'] == 'Non-F2F',
                                    'value'].value_counts()
numbers = counts_list(nonf2f_counts, ['Ethnic Minority', 'White British'])
counts = op_data_temp_eth['value'].value_counts()
totals = counts_list(counts, ['Ethnic Minority', 'White British'])
show_values_on_bars(ax_op_eth, rounded=0,
                    numbers=[j for j in numbers if j!=0])
show_totals(ax_op_eth, totals)
ax_op_eth.set_ylabel('Percentage of ' + title_text)
legend = ax_op_eth.get_legend()
#Get seaborn legend
handles = legend.legend_handles
ax_op_eth.legend(handles, hue_order, title='Appointment Type',
                    bbox_to_anchor=(1,1))
#Proportion testing
if numbers[0] > 0:
    pval_op_eth = propHypothesisTest(numbers[0]/totals[0],
                                    numbers[1]/totals[1],
                                    numbers[0], numbers[1], alpha=0.025)
    #Add bars to show if there is a significant difference
    slide10_text = ('Significant difference' if pval_op_eth < 0.025
            else ' No significant difference')
else:
    slide10_text = 'No Ethnic Minority Data'
label_diff(ax_op_eth, 0, 1, slide10_text, [0,1], [1,1])
plt.savefig('plots/' + title_text + ' Slide 10.png', bbox_inches='tight')

# =======================================================================
############Ethnicity Non-F2F testing, separated by pfmgt_spec  SLIDE 11
spec_unique = op_data['specialty'].unique()
#Make an empty list to store pvalues in
spec_eth_pvals = []
spec_wb_nf2f_perc = []
spec_me_nf2f_perc = []
n_wb_f2f=[]
n_wb_nf2f=[]
n_me_f2f=[]
n_me_nf2f=[]

for spec in spec_unique:
    #Separate by EM
    op_data_temp = (op_data_temp_eth.loc[op_data_temp_eth['specialty']
                                        == spec].copy()
                                        .sort_values(by='Visit'))
    #Make a list of the numbers to include under the percentages
    #could not shorten as sometimes there isn't data for these.
    numbers = [op_data_temp[(op_data_temp['value'] == 'White British')
                            & (op_data_temp['Visit']
                                == 'Non-F2F')].shape[0],
                op_data_temp[(op_data_temp['value'] == 'Ethnic Minority')
                            & (op_data_temp['Visit']
                                == 'Non-F2F')].shape[0]]
    totals = [op_data_temp[(op_data_temp['value']
                            == 'White British')].shape[0],
                op_data_temp[(op_data_temp['value']
                            == 'Ethnic Minority')].shape[0]]
    n_wb_nf2f.append(numbers[0])
    n_me_nf2f.append(numbers[1])
    n_wb_f2f.append(totals[0] - numbers[0])
    n_me_f2f.append(totals[1] - numbers[1])
    #If either of the totals values is zero, set as zero 
    if totals[0] == 0:
        spec_wb_nf2f_perc.append(0)
    else:
        spec_wb_nf2f_perc.append(100*(numbers[0]/totals[0]))
    if totals[1] == 0:
        spec_me_nf2f_perc.append(0)
    else:
        spec_me_nf2f_perc.append(100*(numbers[1]/totals[1]))
    
    #If no non-F2F patients are seen in either or both ethnicity categories, 
    #cannot perform hypothesis test
    if ((all(j > 50 for j in numbers)) or
        ((all(j > 50 for j in totals))and (all(j > 0 for j in numbers)))):
        spec_eth_pvals.append(propHypothesisTest(numbers[0]/totals[0],
                                                    numbers[1]/totals[1],
                                                    numbers[0], numbers[1],
                                                    alpha=0.025))
    else:
        spec_eth_pvals.append(None)
#create dataframe
spec_eth_pvals = pd.DataFrame({'Specialty':spec_unique,
                                'p-value':spec_eth_pvals,
                                'WB Non-F2F Perc':spec_wb_nf2f_perc,
                                'EM Non-F2F Perc':spec_me_nf2f_perc,
                                'WB F2F':n_wb_f2f, 'WB NF2F':n_wb_nf2f,
                                'EM F2F':n_me_f2f, 'EM NF2F':n_me_nf2f})
spec_eth_pvals['WB-EM NonF2F Perc'] = (spec_eth_pvals['WB Non-F2F Perc']
                                        - spec_eth_pvals['EM Non-F2F Perc'])
#Barplot
nf2f_eth_slc_plt = pd.melt(spec_eth_pvals[spec_eth_pvals['p-value']<0.025]
                            .sort_values('WB-EM NonF2F Perc'),
                            id_vars=['Specialty'],
                            value_vars = ['EM Non-F2F Perc',
                                            'WB Non-F2F Perc'])
fig, ax_nf2fspec_eth = plt.subplots(1,1,figsize=(9,4))
sns.barplot(data=nf2f_eth_slc_plt, x='Specialty', hue='variable', y='value',
            ax=ax_nf2fspec_eth, palette=['royalblue','lightskyblue'])
legend = ax_nf2fspec_eth.get_legend()
#Get seaborn legend
handles = legend.legend_handles
ax_nf2fspec_eth.legend(handles, ['Ethnic Minority', 'White British'],
                        bbox_to_anchor=(1,1))
if title_text == 'Outpatients':
    ax_nf2fspec_eth.set_ylabel(tw.fill('Percentage of ' + title_text
                                        + ' seen non-F2F', 40))
else:
    ax_nf2fspec_eth.set_ylabel(tw.fill('Percentage of ' + title_text
                                        + ' carried out non-F2F', 40))
ax_nf2fspec_eth.set_xticks(ax_nf2fspec_eth.get_xticks())
labels = [tw.fill(l, 16) for l in nf2f_eth_slc_plt.Specialty.unique()]
ax_nf2fspec_eth.set_xticklabels(labels = labels)
plt.tight_layout()
plt.savefig('plots/' + title_text + ' Slide 11.png', bbox_inches='tight')

# ====================================================================
#         #Age  SLIDE 12
# ====================================================================
op_data_temp_age = op_data.copy()
#group ages in age bands (-1 not 0 to include babies)
age_list = ['0-19', '20-29', '30-39', '40-49', '50-59', '60-69', '70-79',
            '80+']
op_data_temp_age['value'] = pd.cut(x=op_data_temp_age['Age'],
                                bins=[-1, 19, 29, 39, 49, 59, 69, 79, 120],
                                labels=age_list)
#Make a plot of the proportions
fig1, ax_op_age=plt.subplots(1, 1, figsize = (7,5))
sns.histplot(data=op_data_temp_age.sort_values('value'), x='value',
                hue="Visit", hue_order=hue_order, multiple="fill",
                discrete=True, palette=['lemonchiffon','orange'], alpha=1,
                ax=ax_op_age)
#Make a list of the numbers to include under the percentages
nonf2f_counts = op_data_temp_age.loc[op_data_temp_age['Visit'] == 'Non-F2F',
                                        'value'].value_counts()
num_list = [nonf2f_counts[band] if band in nonf2f_counts.index else 0
            for band in age_list]
counts = op_data_temp_age['value'].value_counts()
tot_list = [counts[band] if band in counts.index
            else 0 for band in age_list]
show_values_on_bars(ax_op_age,rounded=0,
                    numbers = [j for j in num_list if j!=0])
show_totals(ax_op_age, tot_list)
ax_op_age.yaxis.set_major_formatter(PercentFormatter(1))
ax_op_age.set_xlabel('Age')
ax_op_age.set_ylabel('Percentage of ' + title_text)
legend = ax_op_age.get_legend()
# #Get seaborn legend
handles = legend.legend_handles
ax_op_age.legend(handles, hue_order, title='Appointment Type',
                    bbox_to_anchor=(1,1))
plt.savefig('plots/' + title_text + ' Slide 12.png', bbox_inches='tight')

#Hypothesis testing for age bands
#Get all pairs of decile numbers with no repeats
age_pvals = []
for r in itertools.combinations(range(len(age_list)),2):
    pval = propHypothesisTest(num_list[r[0]]/tot_list[r[0]],
                                num_list[r[1]]/tot_list[r[1]],
                                num_list[r[0]], num_list[r[1]])
    #Record NON statistically significant results (as shorter to write)
    if pval > 0.025:
        age_pvals.append([r, pval])

#Make a op_data of the results
age_prop_test = pd.DataFrame(age_pvals, columns=['pairs', 'pval'])
age_prop_test['age1'] = (age_prop_test['pairs']
                            .apply(lambda x: age_list[x[0]]))
age_prop_test['age2'] = (age_prop_test['pairs']
                            .apply(lambda x: age_list[x[1]]))

#write a string of the age pairs that have a significant difference
pairs = (age_prop_test.groupby('age1', as_index=False)['age2']
            .apply(list).values.tolist())
slide12_text = ''
for pair in pairs:
    string = f'{pair[0]} & {pair[1][0]}\n'
    slide12_text += string

# =============================================================================
# #DM01 Queries
# =============================================================================
t2 = time.time()
print(f'First analysis run in {(t2-t1)/60} mins')
print('DM01 query section')
dm01_endo_query = f"""
							SET NOCOUNT ON;
DECLARE				@start DATETIME, 
					@end DATETIME
SET					@start='01-Apr-2021'
SET					@end='{op_end_date}'--Change start and end dates here

--	SECTION ONE *******************************************************************************************

SELECT Proc_code AS op_code, proce AS op_desc, op = 1,
op_type = CASE WHEN Proc_Group = 'OGD' THEN 'Gastro' ELSE Proc_Group END,
op_numb = CASE WHEN Proc_Group = 'OGD' THEN 1 WHEN Proc_Group = 'Colon' THEN 2
ELSE 4 END
INTO                
#procs --drop table #procs
FROM InfoDB.dbo.Diag_treat_proc_codes
WHERE Proc_Group IN ('COLON','OGD','FLEXI','Video')
--add index ON code
CREATE INDEX		op_code_ind 
ON					#procs (op_code)
--57

--SECTION TWO ************************************************************************************************

SELECT ip.pasid, [Referral date] = b.wlist_dttm,
[Attendance date] = ip.admit_dttm, ip.admit_dttm, ip.disch_dttm,
ip.pat_age_on_admit, ip.pat_dob, ip.local_spec,
proc_date = CASE WHEN proc0.op_code IS NOT NULL THEN ip.main_proc_date
			     WHEN proc1.op_code IS NOT NULL THEN ip.proc1_date
			     WHEN proc2.op_code IS NOT NULL THEN ip.proc2_date									  
			     WHEN proc3.op_code IS NOT NULL THEN ip.proc3_date
			     WHEN proc4.op_code IS NOT NULL THEN ip.proc4_date
			     WHEN proc5.op_code IS NOT NULL THEN ip.proc5_date
			     WHEN proc6.op_code IS NOT NULL THEN ip.proc6_date
			     WHEN proc7.op_code IS NOT NULL THEN ip.proc7_date
			     WHEN proc8.op_code IS NOT NULL THEN ip.proc8_date
			     WHEN proc9.op_code IS NOT NULL THEN ip.proc9_date END,  
proc_type = ISNULL(proc0.op_type,
            ISNULL(proc1.op_type,
			ISNULL(proc2.op_type,
			ISNULL(proc3.op_type,
			ISNULL(proc4.op_type,''))))), 
proc_code = ISNULL(proc0.op_code,
			ISNULL(proc1.op_code,
			ISNULL(proc2.op_code,
			ISNULL(proc3.op_code,
			ISNULL(proc4.op_code,''))))),
proc_desc =	ISNULL(proc0.op_desc,
			ISNULL(proc1.op_desc,
			ISNULL(proc2.op_desc,
			ISNULL(proc3.op_desc,
			ISNULL(proc4.op_desc,''))))),
n_procs = ISNULL(proc0.op,0)
		 +ISNULL(proc1.op,0)
		 +ISNULL(proc2.op,0)
		 +ISNULL(proc3.op,0)	
		 +ISNULL(proc4.op,0)
		 +ISNULL(proc5.op,0)
		 +ISNULL(proc6.op,0)
		 +ISNULL(proc7.op,0)
		 +ISNULL(proc8.op,0)
		 +ISNULL(proc9.op,0), --number of procs FROM list
proc_num = ISNULL(proc0.op_numb,0)|
		   ISNULL(proc1.op_numb,0)|
		   ISNULL(proc2.op_numb,0)|
		   ISNULL(proc3.op_numb,0)|
		   ISNULL(proc4.op_numb,0)
		   +ISNULL(proc5.op_numb,0)|
		   ISNULL(proc6.op_numb,0)|
		   ISNULL(proc7.op_numb,0)|
		   ISNULL(proc8.op_numb,0)|
		   ISNULL(proc9.op_numb,0), 
--proc numbs FROM list
main_proc =	ISNULL(ip.main_proc,''),
ip.main_proc_date, 
Main_proc_desc = ISNULL(proc0.op_desc,''), 
proc1 = ISNULL(ip.proc1,''),
proc2 = ISNULL(ip.proc2,''),
proc3 = ISNULL(ip.proc3,''), 
proc4 = ISNULL(ip.proc4,''), 
proc5 = ISNULL(ip.proc5,''),
proc6 = ISNULL(ip.proc6,''), 
proc7 = ISNULL(ip.proc7,''), 
proc8 = ISNULL(ip.proc8,''), 
proc9 = ISNULL(ip.proc9,''),
ip.prvsp_refno, 
ip.ffce_yn, 
-- new section --
pcg_code = InfoDB.dbo.fn_ccg(ip.registered_commissioner,
ip.resident_commissioner),
--
ip.reg_practice_code,
admet = ip.admet_flag + CASE WHEN ip.admet='13' THEN ' PL' 
							 WHEN ip.admet IN ('11','12') THEN ' WL' ELSE '' END
INTO #pats
FROM PimsMarts.dbo.inpatients ip
LEFT JOIN #procs AS proc0 
ON proc0.op_code = ip.main_proc	
LEFT JOIN #procs AS proc1 
ON proc1.op_code = ip.proc1	
LEFT JOIN #procs AS proc2 
ON proc2.op_code = ip.proc2	
LEFT JOIN #procs AS proc3 
ON proc3.op_code = ip.proc3
LEFT JOIN #procs AS proc4 
ON proc4.op_code = ip.proc4	
LEFT JOIN #procs AS proc5 
ON proc5.op_code = ip.proc5
LEFT JOIN #procs AS proc6 
ON proc6.op_code = ip.proc6	
LEFT JOIN #procs AS proc7 
ON proc7.op_code = ip.proc7
LEFT JOIN #procs AS proc8 
ON proc8.op_code = ip.proc8	
LEFT JOIN #procs AS proc9 
ON proc9.op_code = ip.proc9
LEFT JOIN PiMSMarts.dbo.waiting_lists_ipdc_additions b
ON ip.wlist_refno = b.wlist_refno
WHERE ip.fce_yn = 'y'
AND ((ip.main_proc_date BETWEEN @start AND @end AND proc0.op_code IS NOT NULL)
	OR (ip.proc1_date BETWEEN @start AND @end AND proc1.op_code IS NOT NULL)
	OR (ip.proc2_date BETWEEN @start AND @end AND proc2.op_code IS NOT NULL)
	OR (ip.proc3_date BETWEEN @start AND @end AND proc3.op_code IS NOT NULL)
	OR (ip.proc4_date BETWEEN @start AND @end AND proc4.op_code IS NOT NULL)
	OR (ip.proc5_date BETWEEN @start AND @end AND proc5.op_code IS NOT NULL)
	OR (ip.proc6_date BETWEEN @start AND @end AND proc6.op_code IS NOT NULL)
	OR (ip.proc7_date BETWEEN @start AND @end AND proc7.op_code IS NOT NULL)
	OR (ip.proc8_date BETWEEN @start AND @end AND proc8.op_code IS NOT NULL)
	OR (ip.proc9_date BETWEEN @start AND @end AND proc9.op_code IS NOT NULL))
AND	 provider = 'RK900'
--IN ('RK900','89999','89006','89997','NT100','NT200','NT300','NT400','NTY00')	--just PHNT hospitals
AND patcl IN ('1','2','8')	-- 1-ordinary, 2-day cases, 8-others (excluding regular admissions)
AND fce_adcat <> '02'		-- remove privates - were included AND flagged 
AND ip.admet <> '13' ---exclude planned/surveillance
AND	ip.admet_flag <> 'EM' --- exclude emergency
ORDER BY ip.main_proc_date

-- 641
--SECTION THREE *********************************************************************************************		
--set any 13Q etc to X24  - also make left 3

UPDATE			#pats
SET				pcg_code ='X24'
WHERE			pcg_code ='13Q00'

UPDATE			#pats
SET 			pcg_code = LEFT(pcg_code,3)

--SECTION FOUR ****************************************************************************************

SELECT pasid, [Referral date], [Attendance date], prvsp_refno, pcg_code, admet,
       local_spec, first_proc_date=MIN(proc_date), Gastro=MAX(proc_num&1),
       Colon=MAX((proc_num&2)/2), FlexiSig=MAX((proc_num&4)/4)
INTO #pat_ops
FROM #pats
WHERE proc_date BETWEEN @start AND @end
GROUP BY pasid, prvsp_refno, pcg_code, admet, local_spec, [Referral date],
[Attendance date]

--SECTION FIVE ******************************************************************************************

SELECT proc_mth=CONVERT(char(6), first_proc_date, 112), Op='Gastro', pcg_code, 
admet, local_spec, pasid, [Referral date], [Attendance date]
INTO #ip_ops
FROM #pat_ops
WHERE Gastro = 1

UNION ALL
SELECT proc_mth=CONVERT(char(6), first_proc_date, 112), Op='Colon', pcg_code,
admet, local_spec, pasid, [Referral date], [Attendance date]
FROM #pat_ops
WHERE Colon = 1

UNION ALL
SELECT proc_mth=CONVERT(char(6), first_proc_date, 112), Op='FlexiSig', pcg_code,
admet, local_spec, pasid, [Referral date], [Attendance date]
FROM #pat_ops
WHERE FlexiSig = 1

--SECTION SIX *******************************************************************************************

SELECT pasid, [Referral date]=ref_recvd_dttm, [Attendance date]=arrived_dttm,
local_spec, start_dttm, year_month,
pcg_code=InfoDB.dbo.fn_ccg(op.registered_commissioner,
op.resident_commissioner),
op.reg_practice_code, op.schdl_refno, op.comments,     
proc_type =	ISNULL(proc0.op_type,
			ISNULL(proc1.op_type,
			ISNULL(proc2.op_type,
			ISNULL(proc3.op_type,
			ISNULL(proc4.op_type,''))))), 
proc_code =	ISNULL(proc0.op_code,
			ISNULL(proc1.op_code,
			ISNULL(proc2.op_code,
			ISNULL(proc3.op_code,
			ISNULL(proc4.op_code,''))))),
proc_desc =	ISNULL(proc0.op_desc,
			ISNULL(proc1.op_desc,
			ISNULL(proc2.op_desc,
			ISNULL(proc3.op_desc,
			ISNULL(proc4.op_desc,''))))),
n_procs = ISNULL(proc0.op,0)
		  +ISNULL(proc1.op,0)
		  +ISNULL(proc2.op,0)
		  +ISNULL(proc3.op,0)
		  +ISNULL(proc4.op,0)
		  +ISNULL(proc5.op,0), --number of procs FROM list
proc_num = ISNULL(proc0.op_numb,0)|
		   ISNULL(proc1.op_numb,0)|
		   ISNULL(proc2.op_numb,0)|
		   ISNULL(proc3.op_numb,0)|
		   ISNULL(proc4.op_numb,0)|
		   ISNULL(proc5.op_numb,0), --proc numbs FROM list
main_proc =	ISNULL(op.main_proc,''), op.main_proc_date, 
Main_proc_desc = ISNULL(proc0.op_desc,''), 
proc1 = ISNULL(op.proc1,''),
proc2 = ISNULL(op.proc2,''), 
proc3 = ISNULL(op.proc3,''), 
proc4 = ISNULL(op.proc4,''), 
proc5 = ISNULL(op.proc5,''),
coded_yn = CASE	WHEN main_proc IS NULL
AND proc_concat IS NULL THEN 'n' ELSE 'y' END,
IP_yn = CASE WHEN main_proc = 'IP' THEN 'y' ELSE 'n' end,
com_proc = CASE WHEN op.comments LIKE '%flexi%' THEN 'Flexi'
WHEN (op.comments LIKE '%double%' AND op.comments LIKE '%ender%')
OR (op.comments LIKE '%colon%' AND op.comments LIKE '%OGD%') THEN 'Double'
WHEN op.comments LIKE '%colon%' THEN 'Colon'
WHEN op.comments LIKE '%OGD%' THEN 'Gastro'
WHEN op.comments LIKE '%ERCP%' OR op.comments LIKE 'PEG%' THEN 'Other'
ELSE '???' END
INTO #opats
FROM PiMSMarts.dbo.outpatients op
LEFT JOIN #procs AS proc0 
ON proc0.op_code = LEFT(op.main_proc,4)
LEFT JOIN #procs AS proc1 
ON proc1.op_code = LEFT(op.proc1,4)	
LEFT JOIN #procs AS proc2 
ON proc2.op_code = LEFT(op.proc2,4)
LEFT JOIN #procs AS proc3 
ON proc3.op_code = LEFT(op.proc3,4)
LEFT JOIN #procs AS proc4 
ON proc4.op_code = LEFT(op.proc4,4)
LEFT JOIN #procs AS proc5 
ON proc5.op_code = LEFT(op.proc5,4)
WHERE local_spec = '64'
--AND start_dttm BETWEEN '01-Dec-2020' AND '31-Dec-2020 23:59:59'
AND start_dttm BETWEEN @start AND @end
AND attnd = '5'  /****ATT*****/
AND sctyp = 'OTPAT'
AND conlt_yn = 'Y'
AND adcat <> '02'		-- exclude privates
AND	(main_proc <> '499BS' OR main_proc IS NULL)	-- excludes national bowel screening pts	
AND	visit IN ('1','2')
AND op.wlist_refno IS NOT NULL -- exclude any emergency visits and second half of 2-part appts
ORDER BY year_month

--239

--SELECT DISTINCT pcg_code FROM #opats

UPDATE			#opats
SET 			pcg_code ='X24'
WHERE			pcg_code ='13Q00'

UPDATE			#opats
SET 			pcg_code = LEFT(pcg_code,3)

--SECTION EIGHT ****************************************************************************

SELECT pasid, [Referral date], [Attendance date], schdl_refno, start_dttm,
year_month, pcg_code, local_spec, Gastro=MAX(proc_num&1),
Colon=MAX((proc_num&2)/2), FlexiSig=MAX((proc_num&4)/4),
Other=SUM(CASE WHEN proc_num = 0 AND coded_yn = 'y' AND IP_yn = 'n' THEN 1 
ELSE 0 end),
Uncoded = SUM(CASE WHEN coded_yn = 'n' THEN 1 ELSE 0 END),
IP  = SUM(CASE WHEN IP_yn = 'y' THEN 1 ELSE 0 end)	
INTO #opat_ops
FROM #opats
GROUP BY pasid, schdl_refno, start_dttm, year_month, pcg_code, local_spec,
[Referral date], [Attendance date]

-- SECTION NINE *********************************************************************************************

SELECT pasid, [Referral date], [Attendance date], schdl_refno, start_dttm,
year_month, pcg_code, local_spec,
type =	CASE WHEN coded_yn='n' THEN 'Uncoded'
             WHEN IP_yn = 'y' THEN 'IP' ELSE '?' END,
Gastro = SUM( CASE com_proc	WHEN 'Gastro' THEN 1 ELSE 0 END),
Colon = SUM(CASE com_proc WHEN 'Colon' THEN 1 ELSE 0 END),
D_Ender = SUM(CASE com_proc WHEN 'Double' THEN 1 ELSE 0 END),
FlexiSig = SUM(CASE com_proc WHEN 'Flexi' THEN 1 ELSE 0 END),
Other   = SUM(CASE com_proc WHEN 'Other' THEN 1 ELSE 0 END),
Unknown = SUM(CASE com_proc WHEN '???' THEN 1 ELSE 0 END)
INTO #unc_ops
FROM #opats
WHERE coded_yn='n' OR IP_yn = 'y'
group by pasid, [Referral date], [Attendance date], schdl_refno, start_dttm,
year_month, pcg_code, local_spec, 
CASE WHEN coded_yn='n' THEN 'Uncoded' WHEN IP_yn = 'y' THEN 'IP' ELSE '?' END

-- SECTION TEN *****************************************************************************

SELECT proc_mth=year_month, Op='Gastro', pcg_code, admet='OP', local_spec,
pasid, [Referral date], [Attendance date]	
INTO #op_ops
FROM #opat_ops
WHERE Gastro = 1 

UNION ALL
SELECT proc_mth=year_month, Op='Colon', pcg_code, admet='OP', local_spec, pasid,
[Referral date], [Attendance date]
FROM #opat_ops
WHERE Colon = 1

UNION ALL
SELECT proc_mth=year_month, Op='FlexiSig', pcg_code, admet='OP', local_spec,
pasid, [Referral date], [Attendance date]
FROM #opat_ops
WHERE FlexiSig = 1

--SECTION ELEVEN *******************************************************************************
--OP IPs AND uncoded

SELECT proc_mth=year_month, Op='Gastro', pcg_code, admet='OP', local_spec,
pasid, [Referral date], [Attendance date]
INTO #op_unc_ops
FROM #unc_ops
WHERE Gastro = 1 OR D_Ender = 1

UNION ALL
SELECT proc_mth=year_month, Op='Colon', pcg_code, admet='OP', local_spec, pasid,
[Referral date], [Attendance date]
FROM #unc_ops
WHERE Colon = 1 OR D_Ender = 1

UNION ALL
SELECT proc_mth=year_month, Op='FlexiSig', pcg_code, admet='OP', local_spec,
pasid, [Referral date], [Attendance date]
FROM #unc_ops
WHERE FlexiSig = 1

--SECTION TWELVE *****************************************************************************************	

--IPDC Endos
SELECT pasid, [Referral date], [Attendance date], ptype='IPDC', refno=prvsp_refno,
pcg_code, admet, local_spec, proc_date=first_proc_date, op='Gastro'
INTO #all_ops
FROM #pat_ops
WHERE Gastro=1
UNION ALL
SELECT pasid, [Referral date], [Attendance date], ptype='IPDC', refno=prvsp_refno,
pcg_code, admet, local_spec, proc_date=first_proc_date, op='Colon'
FROM #pat_ops
WHERE Colon=1
UNION ALL
SELECT pasid, [Referral date], [Attendance date], ptype='IPDC', refno=prvsp_refno,
pcg_code, admet, local_spec, proc_date=first_proc_date, op='FlexiSig'
FROM #pat_ops
WHERE FlexiSig=1
UNION ALL
--Outpatient Endos
SELECT pasid, [Referral date], [Attendance date], ptype='OP', refno=schdl_refno,
pcg_code, admet='EL WL', local_spec, proc_date=start_dttm, op='Gastro'
FROM #opat_ops
WHERE Gastro=1
UNION ALL
SELECT pasid, [Referral date], [Attendance date], ptype='OP', refno=schdl_refno,
pcg_code, admet='EL WL', local_spec, proc_date=start_dttm,
op='Colon'
FROM #opat_ops
WHERE Colon=1
UNION ALL
SELECT pasid, [Referral date], [Attendance date], ptype='OP', refno=schdl_refno,
pcg_code, admet='EL WL', local_spec, proc_date=start_dttm, op='FlexiSig'
FROM #opat_ops
WHERE FlexiSig=1
UNION ALL
--Uncoded Outpatients
SELECT pasid, [Referral date], [Attendance date], ptype='OP', refno=schdl_refno,
pcg_code, admet='EL WL', local_spec, proc_date=start_dttm, op='Gastro'
FROM #unc_ops
WHERE Gastro=1 OR D_Ender=1
UNION ALL
SELECT pasid, [Referral date], [Attendance date], ptype='OP', refno=schdl_refno,
pcg_code, admet='EL WL', local_spec, proc_date=start_dttm,
op='Colon'
FROM #unc_ops
WHERE Colon=1 OR D_Ender=1
UNION ALL
SELECT pasid, [Referral date], [Attendance date], ptype='OP', refno=schdl_refno,
pcg_code, admet='EL WL', local_spec, proc_date=start_dttm,
op='FlexiSig'
FROM #unc_ops
WHERE FlexiSig=1 

SELECT #all_ops.pasid, #all_ops.[Referral date], #all_ops.[Attendance date],
[Eth].[description] as Ethnicity, imd.[IndexValue]
FROM #all_ops
LEFT JOIN PiMSMarts.dbo.patients pats
ON #all_ops.pasid = pats.pasid
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth 
ON Eth.identifier = pats.ethgr
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
ON pats.pat_pcode = imd.PostcodeFormatted

WHERE [Referral date] IS NOT NULL
AND [Attendance date] IS NOT NULL
AND imd.EndDate IS NULL
ORDER BY [Attendance date]"""
print('Endo query running...')
endo_DM01 = pd.read_sql(dm01_endo_query, sdmart_engine)
print('Endo query complete')

print('CRIS query running...')
dm01_cris_query = f"""
SELECT pasid = act.Pasid, [Referral date]=ISNULL(last_wl_date,Request_date),
[Attendance date] = Event_Date ,Ethnicity = Eth.[description],
imd.IndexValue
FROM [InfoDB].[dbo].CRIS_Activity_Data act

--join the waiting list to find the clock start date
LEFT JOIN (SELECT DISTINCT wldte.event_key, MAX(wldte.wlist_date) AS last_wl_date
		   FROM [InfoDB].[dbo].[CRIS_Weekly_Waiting_List_Data] wldte
		   --inner join #activity on #activity.Event_key = wldte.event_key
		   GROUP BY wldte.event_key) wld
ON wld.event_key = act.Event_key

--join the CRIS exam codes look-up to get a list of the ECG procedures to include
LEFT JOIN [InfoDB].[dbo].[CRIS_exam_codes] codes
ON codes.Code = act.Examination
--join ethnicity and IMD tables
LEFT JOIN PiMSMarts.dbo.patients pats
ON act.pasid = pats.pasid
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth 
ON Eth.identifier = pats.ethgr
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
ON pats.pat_pcode = imd.PostcodeFormatted

WHERE Planned_yn = 'N'
AND Urg NOT IN (8,9)--<> 9 -- Removing urgency = planned and treatment 
AND Pat_Type NOT IN ('AE','IP') --Discount already admitted patients and ED patients
AND (act.Modality IN ('CT','Ultrasound','MRI','DEXA') --'Radiology' removed as not in DM01 guidance
	or act.Examination IN ('FBACE','FBAEN','FBAIY','FSBEN','FBASM','FBASW') -- Barium enema and barium swallow, as stated in DM01 guidance
	or codes.CD_Echo_yn = 'y') -- SELECT only echocardiograms from cardiology. EPS are inpatients and included in DMO1 query
AND Event_Date BETWEEN '01-Apr-2021' AND '{op_end_date}'
AND imd.EndDate IS NULL
ORDER BY CAST(Event_Date AS date) DESC"""
cris_DM01 = pd.read_sql(dm01_cris_query, sdmart_engine)
print('CRIS query complete')

print('Audiology query running...')
dm01_audio_query = f"""
SELECT  pasid = audio.pasid,
[Referral date] = CASE WHEN opact.visit = '2' THEN op_add.wlist_dttm -- if it's a follow-up, look at when they started on the WL
					   WHEN opact.schdl_refno IS NOT NULL THEN opact.ref_recvd_dttm --if it's a new appt, use the referral time
					   WHEN noncharge.CalenderID IS NOT NULL THEN CAST(audio.start_dttm AS date) --walk in
					   END,
[Attendance date] = audio.start_dttm, Ethnicity = Eth.[description], imd.IndexValue
FROM InfoDB.dbo.Audiology_assessments_202324 audio
LEFT JOIN PimsMarts.dbo.outpatients opact
ON audio.schdl_refno = opact.schdl_refno AND audio.pasid = opact.pasid  --join on both as without pasid, joins different patients!
LEFT JOIN InfoDB.dbo.auditbase_activity_noncharge noncharge
ON audio.schdl_refno = noncharge.CalenderID
LEFT JOIN [PiMSMarts].[dbo].[waiting_lists_op_additions] op_add
ON opact.wlist_refno = op_add.wlist_refno
--join ethnicity and IMD tables
LEFT JOIN PiMSMarts.dbo.patients pats
ON audio.pasid = pats.pasid
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth 
ON Eth.identifier = pats.ethgr
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
ON pats.pat_pcode = imd.PostcodeFormatted
WHERE (audio.Activity_group='ASSESS' OR audio.Activity_group1='ASSESS')
AND NOT (audio.Activity_group = 'Fittings' OR audio.Activity_group1 ='Fittings') --Remove fit and assess appts as these are not DM01
AND audio.start_dttm BETWEEN '01-Apr-2023' AND '{op_end_date}'
AND imd.EndDate IS NULL

UNION ALL
SELECT  pasid = audio.pasid,
[Referral date] = CASE WHEN opact.visit = '2' THEN op_add.wlist_dttm -- if it's a follow-up, look at when they started on the WL
					   WHEN opact.schdl_refno IS NOT NULL THEN opact.ref_recvd_dttm --if it's a new appt, use the referral time
					   WHEN noncharge.CalenderID IS NOT NULL THEN CAST(audio.start_dttm AS date) --walk in
					   END,
[Attendance date] = audio.start_dttm, Ethnicity = Eth.[description], imd.IndexValue
FROM InfoDB.dbo.Audiology_assessments_202223 audio
LEFT JOIN PimsMarts.dbo.outpatients opact
ON audio.schdl_refno = opact.schdl_refno AND audio.pasid = opact.pasid  --join on both as without pasid, joins different patients!
LEFT JOIN InfoDB.dbo.auditbase_activity_noncharge noncharge
ON audio.schdl_refno = noncharge.CalenderID
LEFT JOIN [PiMSMarts].[dbo].[waiting_lists_op_additions] op_add
ON opact.wlist_refno = op_add.wlist_refno
--join ethnicity and IMD tables
LEFT JOIN PiMSMarts.dbo.patients pats
ON audio.pasid = pats.pasid
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth 
ON Eth.identifier = pats.ethgr
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
ON pats.pat_pcode = imd.PostcodeFormatted
WHERE (audio.Activity_group ='ASSESS' OR audio.Activity_group1 ='ASSESS')
AND NOT (audio.Activity_group = 'Fittings' OR audio.Activity_group1 ='Fittings') --Remove fit and assess appts as these are not DM01
AND imd.EndDate IS NULL

UNION ALL
SELECT  pasid = audio.pasid,
[Referral date] = CASE WHEN opact.visit = '2' THEN op_add.wlist_dttm -- if it's a follow-up, look at when they started on the WL
					   WHEN opact.schdl_refno IS NOT NULL THEN opact.ref_recvd_dttm --if it's a new appt, use the referral time
					   WHEN noncharge.CalenderID IS NOT NULL THEN CAST(audio.start_dttm AS date) --walk in
					   END,
[Attendance date] = audio.start_dttm, Ethnicity = Eth.[description], imd.IndexValue
from InfoDB.dbo.Audiology_assessments_202122 audio
LEFT JOIN PimsMarts.dbo.outpatients opact
ON audio.schdl_refno = opact.schdl_refno AND audio.pasid = opact.pasid  --join on both as without pasid, joins different patients!
LEFT JOIN InfoDB.dbo.auditbase_activity_noncharge noncharge
ON audio.schdl_refno = noncharge.CalenderID
LEFT JOIN [PiMSMarts].[dbo].[waiting_lists_op_additions] op_add
ON opact.wlist_refno = op_add.wlist_refno
--join ethnicity and IMD tables
LEFT JOIN PiMSMarts.dbo.patients pats
ON audio.pasid = pats.pasid
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth 
ON Eth.identifier = pats.ethgr
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
ON pats.pat_pcode = imd.PostcodeFormatted
WHERE (audio.Activity_group ='ASSESS' OR audio.Activity_group1 ='ASSESS')
AND NOT (audio.Activity_group = 'Fittings' OR audio.Activity_group1 ='Fittings') --Remove fit and assess appts as these are not DM01
AND imd.EndDate IS NULL
ORDER BY audio.start_dttm
"""
audio_DM01 = pd.read_sql(dm01_audio_query, sdmart_engine)
print('Audiology query complete')

print('Other DM01 query running...')
dm01_other_query = f"""
SET NOCOUNT ON;
DECLARE				@DtmStartDate as datetime
DECLARE				@DtmEndDate as datetime
SET					@DtmStartDate= '01-Apr-2021' -- start of last month  
SET					@DtmEndDate= '{op_end_date}' -- END of last month --

--create neurophys table 
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] = opact.arrived_dttm, opact.year_month,
TType = CASE WHEN opact.wlist_refno IS NOT NULL THEN 'Waiting List'
			 WHEN opact.visit = '2' THEN 'Planned' ELSE 'Unscheduled' END,
proc_group = ISNULL(dproc1.Proc_Group,
			 ISNULL(dproc2.Proc_Group,
			 ISNULL(dproc3.Proc_Group,'Unk'))),
 proc_code = CASE WHEN ISNULL(dproc1.Proc_Group,
					   ISNULL(dproc2.Proc_Group,
					   ISNULL(dproc3.Proc_Group,'Unk')))='NEURO' 
					   AND list_name LIKE '%tech%' THEN 'EMGTECH'
				  WHEN ISNULL(dproc1.Proc_Group,
					   ISNULL(dproc2.Proc_Group,
					   ISNULL(dproc3.Proc_Group,'Unk')))='NEURO' THEN 'EMGCONS'
				  ELSE ISNULL(dproc1.Proc_Code,
					   ISNULL(dproc2.map,
					   ISNULL(dproc3.Proc_Code,'Unk'))) END,
ip = CASE WHEN ip.prcae_refno IS NULL THEN 'N' ELSE 'Y' END,
ccg_code = LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
				opact.resident_commissioner),3)

INTO #ALLNP -- all neurophys
FROM Pimsmarts.dbo.outpatients opact
LEFT JOIN InfoDB.dbo.Diag_proc_codes dproc1  
ON opact.main_proc = dproc1.proc_code
LEFT JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 ldd 
ON opact.schdl_refno=schdlrefno
LEFT JOIN InfoDB.dbo.neurophysiology_tests dproc2 
ON ldd.testonecode=dproc2.test_code
LEFT JOIN PiMSMarts.dbo.waiting_lists_op_additions wlopadd 
ON opact.wlist_refno=wlopadd.wlist_refno
LEFT JOIN InfoDB.dbo.Diag_proc_codes dproc3  
ON wlopadd.intended_proc = dproc3.proc_code 
LEFT JOIN InfoDB.dbo.vw_ipdc_episode_start_pfmgt ip 
ON opact.pasid=ip.pasid AND ffce_yn='Y'
AND InfoDB.dbo.fn_remove_time(opact.start_dttm) 
BETWEEN InfoDB.dbo.fn_remove_time(ip.admit_dttm) 
AND InfoDB.dbo.fn_remove_time(ISNULL(ip.disch_dttm, GETDATE()))
WHERE opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.local_spec IN ('89','40','A1')
AND cancr_dttm IS NULL
AND attnd='5'
AND (opact.archv_flag IS NULL OR opact.archv_flag = 'N')	
AND	opact.wlist_refno IS NOT NULL -- exclude any emergency visits
AND opact.visit = '1' -- only include new, F2F visits

----------------------------------------------------------------------------------
--create sleepies table

SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] = opact.arrived_dttm, opact.start_dttm, 
test_performed = CASE WHEN testonecode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study' END,
attType=attendancedescription, opact.clinic_code, opact.visit,
test_type = CASE WHEN opact.wlist_refno IS NOT NULL THEN 'Waiting List'
				 WHEN opact.visit = '2' THEN 'Planned' ELSE 'Unscheduled' END, 
ccg_code = LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner, 
				opact.resident_commissioner),3)
into #sleepies
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n 
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s 
ON opact.local_spec = s.local_spec
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit = '1' --only include new, F2F visits
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_spec IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00') 
OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.local_spec IN ('89','40','A1')
AND n.testonecode IN ('MSL','PM','PSN','SE')
AND (attendancetypecode = 'OP' OR attendancetypecode IS NULL)
AND opact.wlist_refno IS NOT NULL -- exclude any planned visits

UNION ALL

SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date]=opact.arrived_dttm, opact.start_dttm, 
test_performed = CASE WHEN testtwocode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study' END,
attType=attendancedescription, opact.clinic_code, opact.visit,
test_type=CASE WHEN opact.wlist_refno IS NOT NULL THEN 'Waiting List'
			   WHEN opact.visit = '2' THEN 'Planned' ELSE 'Unscheduled' END, 
ccg_code=LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
			  opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n 
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s 
ON opact.local_spec = s.local_spec
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit IN ('1')--only include new, F2F visits
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_spec IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%' 
	 AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00') 
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.local_spec IN ('89','40','A1')
AND n.testtwocode IN ('MSL','PM','PSN','SE')
AND (attendancetypecode = 'OP' OR attendancetypecode IS NULL)
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits

UNION ALL
SELECT opact.pasid, [Referral date]=opact.ref_recvd_dttm,
[Attendance date]=opact.arrived_dttm, opact.start_dttm, 
test_performed = CASE WHEN testthreecode IN ('MSL','PM','PSN','SE')
				 THEN 'Sleep Study' END,
attType=attendancedescription, opact.clinic_code, opact.visit,
test_type= CASE WHEN opact.wlist_refno IS NOT NULL THEN 'Waiting List'
				WHEN opact.visit = '2' THEN 'Planned' ELSE 'Unscheduled' END, 
ccg_code=LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
			  opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n 
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s 
ON opact.local_spec = s.local_spec
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit IN ('1') --only include new, F2F visits
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_spec IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00')
	 or (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.local_spec IN ('89','40','A1')
AND n.testthreecode IN ('MSL','PM','PSN','SE')
AND (attendancetypecode = 'OP' OR attendancetypecode IS NULL)
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits

UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] =opact.arrived_dttm, opact.start_dttm, 
test_performed = CASE WHEN testfourcode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study' END,
attType=attendancedescription, opact.clinic_code, opact.visit,
test_type = CASE WHEN opact.wlist_refno IS NOT NULL THEN 'Waiting List'
				 WHEN opact.visit = '2' THEN 'Planned' ELSE 'Unscheduled' END, 
ccg_code = LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n 
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s 
ON opact.local_spec = s.local_spec
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND	opact.visit IN ('1') --only include new, F2F visits
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_spec IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%'
	 AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00')
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.local_spec IN ('89','40','A1')
AND n.testfourcode IN ('MSL','PM','PSN','SE')
AND (attendancetypecode = 'OP' OR attendancetypecode IS NULL)
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits

UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] = opact.arrived_dttm, opact.start_dttm, 
test_performed = CASE WHEN testfivecode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study' END,
attType=attendancedescription, opact.clinic_code, opact.visit,
test_type = CASE WHEN opact.wlist_refno IS NOT NULL THEN 'Waiting List'
				 WHEN opact.visit = '2' THEN 'Planned' ELSE 'Unscheduled' END, 
ccg_code = LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
				opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n 
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s 
ON opact.local_spec = s.local_spec
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit IN ('1')--only include new, F2F visits
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_spec IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%'
	 AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00') OR (opact.provider ='XXXXX' 
	 AND opact.clinic_code ='AC-S'))
AND opact.local_spec IN ('89','40','A1')
AND n.testfivecode IN ('MSL','PM','PSN','SE')
AND (attendancetypecode = 'OP' OR attendancetypecode IS NULL)
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits

UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] =opact.arrived_dttm, opact.start_dttm, 
test_performed = CASE WHEN testsixcode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study' END,
attType=attendancedescription, opact.clinic_code, opact.visit,
test_type = CASE WHEN opact.wlist_refno IS NOT NULL THEN 'Waiting List'
				 WHEN opact.visit = '2' THEN 'Planned' ELSE 'Unscheduled' END, 
ccg_code=LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n 
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s 
ON opact.local_spec = s.local_spec
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit IN ('1')
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_spec IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00')
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.local_spec IN ('89','40','A1')
AND n.testsixcode IN ('MSL','PM','PSN','SE')
AND (attendancetypecode = 'OP' OR attendancetypecode IS NULL)
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits

UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] =opact.arrived_dttm, opact.start_dttm, 
CASE WHEN testonecode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study'
	 ELSE testonedescription END AS test_performed, 
attendancedescription, opact.clinic_code, opact.visit,
CASE WHEN i.admet IN ('11','12') THEN 'Waiting List' 
	 WHEN i.admet = '13' THEN 'Planned'
     WHEN i.admet not IN ('11','12','13') THEN 'Unscheduled' ELSE 'error' END, 
ccg_code = LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
		        opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n 
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s 
ON opact.local_spec = s.local_spec
LEFT OUTER JOIN PimsMarts.dbo.inpatients i 
ON opact.patnt_refno = i.patnt_refno
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit IN ('1','2') -- does this need to exclude FU visits(2)?
AND i.admet IN ('11','12') -- only include elective procedures that are not 'planned'
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_code IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND	(opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00')
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.clinic_code NOT IN ('DOYLEXTV','DOYLEXCB','SKAD1')
AND opact.local_cons_code NOT IN ('FDP','DP')
AND	opact.local_spec IN ('89','40','A1')
AND n.testonecode IN ('MSL','PM','PSN','SE')
AND (n.attendancetypecode = 'IP' OR n.attendancetypecode = 'DC')
AND opact.start_dttm BETWEEN i.admit_dttm AND i.disch_dttm
AND i.ffce_yn = 'Y'
AND	opact.wlist_refno IS NOT NULL -- exclude any emergency visits


UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] = opact.arrived_dttm, opact.start_dttm, 
CASE WHEN testTWOcode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study'
	 ELSE testTWOdescription END as test_performed, 
attendancedescription, opact.clinic_code, opact.visit,
CASE WHEN i.admet IN ('11','12') THEN 'Waiting List' 
	 WHEN i.admet = '13' THEN 'Planned'
	 WHEN i.admet not IN ('11','12','13') THEN 'Unscheduled' ELSE 'error' END, 
ccg_code= LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
			   opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n 
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s 
ON opact.local_spec = s.local_spec
LEFT OUTER JOIN PimsMarts.dbo.inpatients i 
ON opact.patnt_refno = i.patnt_refno
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit IN ('1','2')
AND i.admet IN ('11','12') -- only include elective procedures that are not 'planned'
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_code IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00')
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.clinic_code NOT IN ('DOYLEXTV','DOYLEXCB','SKAD1')
AND opact.local_cons_code NOT IN ('FDP','DP')
AND opact.local_spec IN ('89','40','A1')
AND n.testTWOcode IN ('MSL','PM','PSN','SE')
AND (n.attendancetypecode = 'IP' OR n.attendancetypecode = 'DC')
AND opact.start_dttm BETWEEN i.admit_dttm AND i.disch_dttm
AND i.ffce_yn = 'Y'
AND	opact.wlist_refno IS NOT NULL -- exclude any emergency visits

UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] = opact.arrived_dttm, opact.start_dttm, 
CASE WHEN testTHREEcode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study'
	 ELSE testTHREEdescription END as test_performed, 
attendancedescription, opact.clinic_code, opact.visit,
CASE WHEN i.admet IN ('11','12') THEN 'Waiting List'
	 WHEN i.admet = '13' THEN 'Planned'
	 WHEN i.admet NOT IN ('11','12','13') THEN 'Unscheduled' ELSE 'error' END, 
ccg_code = LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n 
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s 
ON opact.local_spec = s.local_spec
LEFT OUTER JOIN PimsMarts.dbo.inpatients i 
ON opact.patnt_refno = i.patnt_refno
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND	opact.visit IN ('1','2')
and i.admet IN ('11','12') -- only include elective procedures that are not 'planned'
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_code IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND	(opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00') 
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.clinic_code NOT IN ('DOYLEXTV','DOYLEXCB','SKAD1')
AND opact.local_cons_code NOT IN ('FDP','DP')
AND opact.local_spec IN ('89','40','A1')
AND n.testTHREEcode IN ('MSL','PM','PSN','SE')
AND (n.attendancetypecode = 'IP' OR n.attendancetypecode = 'DC')
AND opact.start_dttm BETWEEN i.admit_dttm AND i.disch_dttm
AND i.ffce_yn = 'Y'
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits


UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] = opact.arrived_dttm, opact.start_dttm, 
CASE WHEN testFOURcode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study' ELSE testFOURdescription END as test_performed, 
attendancedescription, opact.clinic_code, opact.visit,
CASE WHEN i.admet IN ('11','12') THEN 'Waiting List' WHEN i.admet = '13' THEN 'Planned'
	 WHEN i.admet not IN ('11','12','13') THEN 'Unscheduled' ELSE 'error' END, 
ccg_code = LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
				opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s
ON opact.local_spec = s.local_spec
LEFT OUTER JOIN PimsMarts.dbo.inpatients i
ON opact.patnt_refno = i.patnt_refno
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit IN ('1','2')
AND i.admet IN ('11','12') -- only include elective procedures that are not 'planned'
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND	opact.session_code IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00')
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.clinic_code NOT IN ('DOYLEXTV','DOYLEXCB','SKAD1')
AND opact.local_cons_code NOT IN ('FDP','DP')
AND opact.local_spec IN ('89','40','A1')
AND n.testFOURcode IN ('MSL','PM','PSN','SE')
AND (n.attendancetypecode = 'IP' OR n.attendancetypecode = 'DC')
AND opact.start_dttm BETWEEN i.admit_dttm AND i.disch_dttm
AND i.ffce_yn = 'Y'
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits

UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] = opact.arrived_dttm, opact.start_dttm, 
CASE WHEN testFIVEcode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study' ELSE testFIVEdescription END AS test_performed, 
attendancedescription, opact.clinic_code, opact.visit,
CASE WHEN i.admet IN ('11','12') THEN 'Waiting List' WHEN i.admet = '13' THEN 'Planned'
	 WHEN i.admet not IN ('11','12','13') THEN 'Unscheduled' ELSE 'error' END, 
ccg_code=LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
		      opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN	PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN	PimsMarts.dbo.cset_specialties s
ON opact.local_spec = s.local_spec
LEFT OUTER JOIN	PimsMarts.dbo.inpatients i
ON opact.patnt_refno = i.patnt_refno
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit IN ('1','2')
AND i.admet in ('11','12') -- only include elective procedures that are not 'planned'
AND	opact.cancr_dttm IS NULL
AND	opact.attnd IN ('5')
AND	opact.session_code IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND	(opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00')
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.clinic_code NOT IN ('DOYLEXTV','DOYLEXCB','SKAD1')
AND opact.local_cons_code NOT IN ('FDP','DP')
AND opact.local_spec IN ('89','40','A1')
AND n.testFIVEcode IN ('MSL','PM','PSN','SE')
AND (n.attendancetypecode = 'IP' OR n.attendancetypecode = 'DC')
AND opact.start_dttm BETWEEN i.admit_dttm AND i.disch_dttm
AND i.ffce_yn = 'Y'
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits

UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] = opact.arrived_dttm, opact.start_dttm, 
CASE WHEN testSIXcode IN ('MSL','PM','PSN','SE') THEN 'Sleep Study' ELSE testSIXdescription END as test_performed, 
attendancedescription, opact.clinic_code, opact.visit,
CASE WHEN i.admet IN ('11','12') THEN 'Waiting List' WHEN i.admet = '13' THEN 'Planned'
	 WHEN i.admet NOT IN ('11','12','13') THEN 'Unscheduled' ELSE 'error' END, 
ccg_code=LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
			  opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n
ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s
ON opact.local_spec = s.local_spec
LEFT OUTER JOIN PimsMarts.dbo.inpatients i
ON opact.patnt_refno = i.patnt_refno
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit IN ('1','2')
AND i.admet IN ('11','12') -- only include elective procedures that are not 'planned'
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND opact.session_code IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00')
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.clinic_code NOT IN ('DOYLEXTV','DOYLEXCB','SKAD1')
AND opact.local_cons_code NOT IN ('FDP','DP')
AND	opact.local_spec IN ('89','40','A1')
AND n.testSIXcode IN ('MSL','PM','PSN','SE')
AND (n.attendancetypecode = 'IP' OR n.attendancetypecode = 'DC')
AND opact.start_dttm BETWEEN i.admit_dttm AND i.disch_dttm
AND i.ffce_yn = 'Y'
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits


UNION ALL
SELECT opact.pasid, [Referral date] = opact.ref_recvd_dttm,
[Attendance date] = opact.arrived_dttm, opact.start_dttm, 
test_performed='Sleep Study-Oximetry', attType='Outpatients',
opact.clinic_code, opact.visit,
test_type = CASE WHEN opact.wlist_refno IS NOT NULL THEN 'Waiting List'
				 WHEN opact.visit = '2' THEN 'Planned' ELSE 'Unscheduled' END, 
ccg_code=LEFT(InfoDB.dbo.fn_ccg(opact.registered_commissioner,
			  opact.resident_commissioner),3)
FROM PimsMarts.dbo.outpatients opact
LEFT OUTER JOIN PiMSMarts.LDD.vw_OutpatientTestsPerformed_126 n ON opact.schdl_refno = n.schdlrefno
LEFT OUTER JOIN PimsMarts.dbo.cset_specialties s ON opact.local_spec = s.local_spec
WHERE opact.sctyp='OTPAT'
AND opact.start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND opact.visit ='1' -- new only 
AND opact.comments LIKE '%oxim%'
AND opact.cancr_dttm IS NULL
AND opact.attnd IN ('5')
AND	opact.session_spec IS NOT NULL
AND (opact.archv_flag IS NULL OR opact.archv_flag ='N')
AND (opact.location_code NOT LIKE '5F1%' AND opact.location_code NOT LIKE '%PCT%')
AND (opact.provider IN('RK900','89006','89999','NT200','NTY00')
	 OR (opact.provider ='XXXXX' AND opact.clinic_code ='AC-S'))
AND opact.wlist_refno IS NOT NULL -- exclude any emergency visits

--create cysto table
SELECT month=CONVERT(CHAR(6),ipact.main_proc_date,112), ipact.nat_spec,
ccg_code=LEFT(InfoDB.dbo.fn_ccg(ipact.registered_commissioner, ipact.resident_commissioner),3),
PCT=CASE WHEN (ipact.mod_yn='Y' or ipact.dha_pcg_code='XMD00') THEN 'MOD'
		 WHEN ipact.dha_pcg_code='TDH00' THEN 'NCA'
		 WHEN ipact.pcg_code IN('49998','NOPCG') THEN '5F100' ELSE ipact.pcg_code END,
Type=CASE WHEN ipact.admet IN ('11','12') THEN 'Waiting List'
		  WHEN ipact.admet IN ('13') THEN 'Planned' ELSE 'Unscheduled' END,
ipact.admet, ipact.ffce_yn, ipact.main_proc, ipact.main_proc_date, ipact.proc_concat,
ipact.hrg_35, ipact.hrg_35_spells, ipact.fce_start_ward, ipact.reg_practice_code,
ipact.pasid, [Referral date] = b.wlist_dttm, [Attendance date] =ipact.admit_dttm
INTO #cysto
FROM PiMSMarts.dbo.inpatients ipact
LEFT JOIN PiMSMarts.dbo.waiting_lists_ipdc_additions b
ON ipact.wlist_refno = b.wlist_refno
LEFT JOIN PiMSMarts.dbo.cset_specialties d
ON ipact.local_spec = d.local_spec

WHERE fce_yn ='Y'
AND ipact.provider IN ('RK900','89999','89006','89997','NT100','NT200','NT300','NT400','NTY00')
AND ipact.main_proc_date BETWEEN @DtmStartDate AND @DtmEndDate
AND ipact.fce_adcat <> '02'  --not private patients
AND ipact.patcl IN ('1','2','8') --Ordinary, DC, NA
AND (ipact.fce_start_ward NOT LIKE '%PCT%' OR ipact.fce_start_ward IS NULL)
AND (ipact.proc_concat LIKE '%M45%' OR ipact.proc_concat LIKE '%M30%' OR ipact.proc_concat LIKE '%M77%'
     OR ipact.main_proc LIKE '%M45%' OR ipact.main_proc LIKE '%M30%' OR ipact.main_proc LIKE '%M77%')
AND ipact.admet IN ('11','12') -- exclude planned and emergency
 ---urodynamics

SELECT start_dttm, sorrfp,
ttype = CASE WHEN wlist_refno IS NOT NULL THEN 'Waiting List'
			 WHEN visit = '2' THEN 'Planned' ELSE 'Unscheduled' END,
attnd, clinic_code, session_code, local_spec, conlt_yn, pfmgt_purch,
wlist_refno, visit, main_proc, proc_concat, comments, pasid,
[Referral date] = ref_recvd_dttm, [Attendance date] =arrived_dttm,
ccg_code=LEFT(InfoDB.dbo.fn_ccg(registered_commissioner, resident_commissioner),3) ,
PCT = CASE WHEN (mod_yn='Y' or dha_pcg_code='XMD00') THEN 'MOD'
  		   WHEN dha_pcg_code='TDH00' THEN 'NCA'
   		   WHEN pcg_code IN('49998','NOPCG') THEN '5F100' ELSE pcg_code END,
reg_practice_code
INTO #urody
FROM PimsMarts.dbo.outpatients
WHERE ((clinic_code LIKE 'URDDF%' AND (main_proc ='M474' OR main_proc ='M07OP' OR main_proc LIKE 'U26%'))
		OR session_code IN ('CA-XL','UROST2X','UROST7A','UROXB7A','UROXB2A','UROFC1','UROFC3'))
AND	sctyp='OTPAT'
AND start_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND visit IN ('1') -- do not include FU (2)
AND cancr_dttm IS NULL
AND attnd IN ('5','9') --include attended or N/A
AND session_code IS NOT NULL
AND pfmgt_purch <>'VPP00'
AND (archv_flag IS NULL OR archv_flag ='N')
AND (location_code NOT LIKE '5F1%' AND location_code NOT LIKE '%PCT%')
AND (provider IN('RK900','89006','89999','NT200','NTY00') OR (provider ='XXXXX' AND clinic_code ='AC-S'))
AND wlist_refno IS NOT NULL --only include patients who have been on a waiting list

----v2.0				update to X24 for all tables prior to extract
UPDATE				#ALLNP
SET					ccg_code = 'X24'
WHERE				Ccg_code IN ('13Q','13H')
UPDATE				#sleepies
SET					ccg_code = 'X24'
WHERE				ccg_code IN ('13Q','13H')
UPDATE				#cysto
SET					ccg_code = 'X24'
WHERE				ccg_code IN ('13Q','13H')
UPDATE				#urody
SET					ccg_code = 'X24'
WHERE				ccg_code IN ('13Q','13H')

--CREATE FINAL TABLE
--create the table for dm01
SELECT PGroup ='Neurophysiology - peripheral neurophysiology', pct=ccg_code,
ttype, pasid, [Referral date], [Attendance date]
INTO #final
FROM #ALLNP
WHERE PROC_GROUP ='NEURO'

UNION ALL
--eps
SELECT PGroup ='Cardiology - Electrophysiology',
pct=LEFT(InfoDB.dbo.fn_ccg(ipact.registered_commissioner,
		 ipact.resident_commissioner),3),
ttype = CASE WHEN ipact.admet ='11' THEN 'Waiting List' ELSE 'Unscheduled' END,
ipact.pasid, [Referral date] = b.wlist_dttm, [Attendance date] = ipact.admit_dttm
FROM pimsmarts.dbo.inpatients ipact
LEFT JOIN PiMSMarts.dbo.waiting_lists_ipdc_additions b
ON ipact.wlist_refno = b.wlist_refno
WHERE (((main_proc = 'K668' AND proc1 = 'G218')
		 OR (proc_concat LIKE '%K668%' AND proc_concat LIKE '%G218%'))
		 OR ((main_proc = 'K582') OR (proc_concat LIKE '%K582%'))
		 OR ((main_proc = 'K572') OR (proc_concat LIKE '%K572%')))
AND fce_yn ='Y'
AND provider = 'RK900'
AND fce_end_dttm BETWEEN @DtmStartDate AND @DtmEndDate
AND (fce_well_baby_yn ='N' OR fce_well_baby_yn IS NULL)
AND ipact.adcat <>'02'
AND patcl IN ('1','2','8')
AND ipact.admet IN ('11','12') --only WL patients. Exclude planned and ED attenders


UNION ALL
SELECT PGroup='Respiratory physiology - sleep studies', pct=ccg_code, ttype = test_type,
pasid, [Referral date], [Attendance date]
FROM #sleepies

UNION ALL
SELECT PGroup='Cystoscopy', pct=ccg_code, ttype =Type, pasid, [Referral date],
[Attendance date]
FROM #cysto

UNION ALL
SELECT PGroup='Urodynamics - pressures & flows', pct=ccg_code, ttype, pasid,
[Referral date], [Attendance date]
FROM #urody

SELECT FF.pasid, [Referral date], [Attendance date], Eth.[description] as Ethnicity,
imd.IndexValue
FROM				#final FF

--join ethnicity and IMD tables
LEFT JOIN PiMSMarts.dbo.patients pats
ON FF.pasid = pats.pasid
LEFT JOIN PiMSMarts.dbo.cset_ethgr Eth 
on Eth.identifier = pats.ethgr
LEFT JOIN [PiMSMarts].[Reference].[vw_IndicesOfMultipleDeprivation2019_DecileByPostcode] imd
on pats.pat_pcode = imd.PostcodeFormatted
where imd.EndDate IS NULL
"""
other_DM01 = pd.read_sql(dm01_other_query, sdmart_engine)
print('Other DM01 query complete')
#All queries run, so close the engine
sdmart_engine.dispose()
t3 = time.time()
print(f'Queries run in {(t3-t2)/60} mins')

# =============================================================================
# #DM01 Analysis Section SLIDE 13
# =============================================================================
print('DM01 analysis running...')
# DM01 Analysis Section
#Concat these all into one df
dm01_full = pd.concat([endo_DM01, audio_DM01, cris_DM01, other_DM01],axis=0)
dm01_full['IndexValue'] = dm01_full['IndexValue'].astype(float)

#Find month and year of attendance
dm01_full['Month'] = pd.DatetimeIndex(dm01_full['Attendance date']).month
dm01_full['Year'] = pd.DatetimeIndex(dm01_full['Attendance date']).year
#Find Length of Wait
dm01_full['LoW'] = (pd.to_datetime(dm01_full['Attendance date'])
                    - pd.to_datetime(dm01_full['Referral date'])).dt.days
#Get rid of any entries with a negative LoW
dm01_full = dm01_full.loc[(dm01_full['LoW'] >= 0)].copy()

#Make subsets of data    
dm01_WB = dm01_full.loc[dm01_full['Ethnicity'] == 'White British'].copy()
dm01_ME = dm01_full.loc[~dm01_full['Ethnicity']
                        .isin(['White British', 'Unknown',
                               'Unwilling to answer'])].copy()
dm01_IMD12 = dm01_full.loc[dm01_full['IndexValue'].isin([1,2])].copy()
dm01_IMD310 = dm01_full.loc[
    dm01_full['IndexValue'].isin([3,4,5,6,7,8,9,10])].copy()
    
#For each category (WB, ME, IMD1&2, IMD3-10), find the median LoW for each month
dfs = [dm01_WB, dm01_ME, dm01_IMD12, dm01_IMD310]
cols = ['White British', 'Ethnic Minority', 'IMD 1&2','IMD 3-10']
#Function for getting medians
def grpDFs(df, col):
    ret_df = df.groupby(['Year','Month'])['LoW'].median().reset_index()
    ret_df['type'] = col
    return ret_df
#loop over each df and append median values to list for plotting
df_to_plot=[]
for df, col in zip(dfs, cols):
    df_to_plot.append(grpDFs(df,col))
df_to_plot = pd.concat(df_to_plot,ignore_index=True)

#Make a column witht he dates for each month for plotting
df_to_plot['Day'] = 1
df_to_plot['Date'] = pd.to_datetime(df_to_plot[['Year','Month','Day']])
#Add in a column to use for clearer plotting
df_to_plot['plot_cat'] = np.where(df_to_plot['type'].isin(['IMD 1&2',
                                                           'IMD 3-10']),
                                                           'IMD', 'Ethnicity')
fig, ax = plt.subplots(2,1,sharex = True)
g = sns.lineplot(data=df_to_plot.loc[df_to_plot['plot_cat'] == 'Ethnicity'],
                 x='Date', y='LoW', hue='type', markers=True,
                 palette=['lightskyblue','royalblue'], ax=ax[0])
h = sns.lineplot(data=df_to_plot.loc[df_to_plot['plot_cat'] == 'IMD'],
                 x='Date', y='LoW', hue='type', markers=True,
                 palette = ['seagreen','lightgreen'], ax=ax[1])
ax[0].set_ylabel('')
ax[1].set_ylabel('')
ax[1].set_xlabel('')
fig.autofmt_xdate()
for plott in [g, h]:
    leg = plott.get_legend()
    #Remove the title
    leg.set_title('')
    #Move the legend
    sns.move_legend(plott, loc='upper left', bbox_to_anchor=(1,1))
fig.supxlabel('')
fig.supylabel('Median Length of Wait (days)')
ax[0].xaxis.set_major_formatter(mdates.DateFormatter('%b %Y'))
#Make yaxes the same
ymin = min(ax[0].get_ylim()[0], ax[1].get_ylim()[0])
ymax = max(ax[0].get_ylim()[1], ax[1].get_ylim()[1])
ax[0].set_ylim(ymin = ymin, ymax = ymax)
ax[1].set_ylim(ymin = ymin, ymax = ymax)
plt.savefig('plots/Slide 13.png', bbox_inches='tight')
t4 = time.time()
print(f'Seccond analysis run in {(t4-t3)/60} mins')
# ====================================================================
#  #Create Powerpoint
# ====================================================================
prs = Presentation()
start = start_date.strftime('%d/%m/%Y')
end = end_date.strftime('%d/%m/%Y')

####SLIDE 1
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = 'Health Inequalities Data'
subtitle.text = version_date

####SLIDE 2
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'RTT LoW – only last 12 months'
img = slide.shapes.add_picture('plots/Slide 2.png',
                               left=Inches(2), top=Inches(1.5))
text_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(6.5),
                                    width=Inches(12), height=Inches(1))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = f'Clock stops from {start} to {end}'
font = run.font
font.size = Pt(24)

####SLIDE 3
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'RTT LoW by SLC - Ethnicity'
img = slide.shapes.add_picture('plots/Slide 3.png',
                               left=Inches(0.5), top=Inches(2))
text_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(6.5),
                                    width=Inches(12), height=Inches(1))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text='Significant difference in median LoW by ethnicity for the SLCs shown'
font = run.font
font.size = Pt(24)

####SLIDE 4
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'RTT LoW by SLC - IMD'
img = slide.shapes.add_picture('plots/Slide 4.png',
                               left=Inches(1.25), top=Inches(1.5))
text_box = slide.shapes.add_textbox(left=Inches(0.75), top=Inches(6.5),
                                    width=Inches(12), height=Inches(1))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = 'Significant difference in median LoW by IMD for the SLCs shown'
font = run.font
font.size = Pt(24)

####SLIDE 5
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'Proportion of patients waiting > 52 weeks'
img = slide.shapes.add_picture('plots/Slide 5.png',
                               left=Inches(2), top=Inches(1.75))
text_box = slide.shapes.add_textbox(left=Inches(0.75), top=Inches(6.5),
                                    width=Inches(12), height=Inches(1))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = f'There is {imd_str} significant difference for ethnicity or IMD'
font = run.font
font.size = Pt(24)

####SLIDE 6
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'Non-F2F Outpatient appointments'
text_box = slide.shapes.add_textbox(left=Inches(0.75), top=Inches(2),
                                    width=Inches(12), height=Inches(1))
tb = text_box.text_frame
tb.text = f'Outpatient appointments from {start} to {end} included'
prg = tb.add_paragraph()
prg.text = 'Cancelled and DNA appointments excluded'

####SLIDE 7
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'IMD – 2 category'
img = slide.shapes.add_picture('plots/Outpatient Appointments Slide 7.png',
                               left=Inches(0.5), top=Inches(2))
text_box = slide.shapes.add_textbox(left=Inches(6), top=Inches(2.25),
                                    width=Inches(2.5), height=Inches(5))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = f'''{slide7_text} in\n percentage of outpatient
 appointments carried\n out non-F2F for IMD 1-2 vs\n IMD 3-10'''
font = run.font
font.size = Pt(24)

####SLIDE 8
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'IMD – 2 category'
img = slide.shapes.add_picture('plots/Outpatient Appointments Slide 8.png',
                               left=Inches(0.5), top=Inches(2))
text_box = slide.shapes.add_textbox(left=Inches(0.75), top=Inches(6),
                                    width=Inches(12), height=Inches(1))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = '''Significant difference in percentage of OP appointments carried
 out non-F2F for the above specialties'''
font = run.font
font.size = Pt(24)

####SLIDE 9
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'IMD – 10 category'
img = slide.shapes.add_picture('plots/Outpatient Appointments Slide 9.png',
                               left=Inches(0.5), top=Inches(2.25),
                               width=Inches(5.5))
text_box = slide.shapes.add_textbox(left=Inches(6), top=Inches(2.25),
                                    width=Inches(3), height=Inches(6))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = f'''Significant difference in\n percentage of outpatient
 appointments carried\n out non-F2F between:\n''' + slide9_text
font = run.font
font.size = Pt(20)

####SLIDE 10
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'Ethnicity'
img = slide.shapes.add_picture('plots/Outpatient Appointments Slide 10.png',
                               left=Inches(0.5), top=Inches(2.25))
text_box = slide.shapes.add_textbox(left=Inches(5.75), top=Inches(2.25),
                                    width=Inches(3.5), height=Inches(5))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = f'''{slide10_text} in\n the proportion of 
non-F2F OP appts for ethnic\n minority and white british\n patients.'''
font = run.font
font.size = Pt(24)

####SLIDE 11
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'Ethnicity by Specialty'
img = slide.shapes.add_picture('plots/Outpatient Appointments Slide 11.png',
                               left=Inches(0.75), top=Inches(2))
text_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(6),
                                    width=Inches(12), height=Inches(1))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = '''Significant difference in percentage of OP appointments carried
 out non-F2F for the above specialties'''
font = run.font
font.size = Pt(24)

####SLIDE 12
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'Age'
img = slide.shapes.add_picture('plots/Outpatient Appointments Slide 12.png',
                               left=Inches(0.5), top=Inches(2.25),
                               width=Inches(5.5))
text_box = slide.shapes.add_textbox(left=Inches(6), top=Inches(2.25),
                                    width=Inches(3), height=Inches(5))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = (f'''Significant difference\n between all age 
categories except:\n''' + slide12_text)
font = run.font
font.size = Pt(24)

####SLIDE 13
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = 'DM01 – excluding LD plot'
img = slide.shapes.add_picture('plots/Slide 13.png',
                               left=Inches(1.75), top=Inches(1.5))
text_box = slide.shapes.add_textbox(left=Inches(1), top=Inches(6),
                                    width=Inches(12), height=Inches(1))
tb = text_box.text_frame
p = tb.paragraphs[0]
run = p.add_run()
run.text = '''Significant difference in DM01 wait times for both pairs of 
characteristics. Sustained decrease in LoW since Jan 2023.'''
font = run.font
font.size = Pt(24)

####SAVE FILE
prs.save(f'outputs/Health Inequalities Data {version_date}.pptx')
print('Presentation created')
t5 = time.time()
print(f'Total run time {(t5-t0)/60} mins')
